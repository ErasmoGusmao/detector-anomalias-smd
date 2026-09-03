"""Identidade visual e blocos de layout da apresentacao do Grupo 12.

Concentra paleta, tipografia e os helpers de posicionamento usados por
``gerar_slides.py``. Nenhuma funcao aqui conhece o conteudo dos slides: a
separacao entre forma (este modulo) e conteudo (``conteudo.py``) segue a mesma
logica de modularizacao adotada no ``src/`` do projeto.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paleta - sobrio academico
# ---------------------------------------------------------------------------
AZUL_HEX = "0F3A52"        # azul-petroleo profundo (cor institucional)
AZUL_CLARO_HEX = "2C6C8F"  # apoio, para textos secundarios sobre fundo claro
AMBAR_HEX = "C8862A"       # destaque, reguas e realces
GRAFITE_HEX = "232A2E"     # texto corrido
CINZA_HEX = "6B7780"       # texto de apoio e rodapes
CINZA_CLARO_HEX = "DCE2E6"  # bordas e linhas de tabela
GELO_HEX = "F7F8FA"        # fundo dos slides
GELO_AZUL_HEX = "EAF1F5"   # fundo de cartoes e zebra de tabela
VERDE_HEX = "2E7D5B"       # atendido / sucesso
VERMELHO_HEX = "A8442C"    # pendente / limitacao
BRANCO_HEX = "FFFFFF"

AZUL = RGBColor.from_string(AZUL_HEX)
AZUL_CLARO = RGBColor.from_string(AZUL_CLARO_HEX)
AMBAR = RGBColor.from_string(AMBAR_HEX)
GRAFITE = RGBColor.from_string(GRAFITE_HEX)
CINZA = RGBColor.from_string(CINZA_HEX)
CINZA_CLARO = RGBColor.from_string(CINZA_CLARO_HEX)
GELO = RGBColor.from_string(GELO_HEX)
GELO_AZUL = RGBColor.from_string(GELO_AZUL_HEX)
VERDE = RGBColor.from_string(VERDE_HEX)
VERMELHO = RGBColor.from_string(VERMELHO_HEX)
BRANCO = RGBColor.from_string(BRANCO_HEX)

# ---------------------------------------------------------------------------
# Tipografia e grade
# ---------------------------------------------------------------------------
FONTE = "Calibri"
FONTE_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGEM = Inches(0.62)
LARGURA_UTIL = SLIDE_W - 2 * MARGEM
TOPO_CONTEUDO = Inches(1.62)
BASE_CONTEUDO = Inches(6.92)
ALTURA_UTIL = BASE_CONTEUDO - TOPO_CONTEUDO


# ---------------------------------------------------------------------------
# Primitivas
# ---------------------------------------------------------------------------
def fundo(slide, cor: RGBColor) -> None:
    """Pinta o fundo do slide com cor solida."""
    preenchimento = slide.background.fill
    preenchimento.solid()
    preenchimento.fore_color.rgb = cor


def retangulo(
    slide,
    left,
    top,
    width,
    height,
    cor=None,
    borda=None,
    espessura_borda=0.75,
    forma=MSO_SHAPE.RECTANGLE,
):
    """Desenha um retangulo (fundo de cartao, faixa, regua)."""
    shape = slide.shapes.add_shape(forma, left, top, width, height)
    if cor is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor
    if borda is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = borda
        shape.line.width = Pt(espessura_borda)
    shape.shadow.inherit = False
    return shape


def texto(
    slide,
    conteudo: str,
    left,
    top,
    width,
    height,
    tamanho=16,
    cor=GRAFITE,
    negrito=False,
    italico=False,
    alinhamento=PP_ALIGN.LEFT,
    fonte=FONTE,
    espaco_linhas=1.0,
    ancora=MSO_ANCHOR.TOP,
):
    """Insere uma caixa de texto simples de um paragrafo."""
    caixa = slide.shapes.add_textbox(left, top, width, height)
    quadro = caixa.text_frame
    quadro.word_wrap = True
    quadro.vertical_anchor = ancora
    quadro.margin_left = 0
    quadro.margin_right = 0
    quadro.margin_top = 0
    quadro.margin_bottom = 0

    paragrafo = quadro.paragraphs[0]
    paragrafo.alignment = alinhamento
    paragrafo.line_spacing = espaco_linhas
    run = paragrafo.add_run()
    run.text = conteudo
    run.font.size = Pt(tamanho)
    run.font.bold = negrito
    run.font.italic = italico
    run.font.color.rgb = cor
    run.font.name = fonte
    return caixa


def paragrafos(
    slide,
    linhas,
    left,
    top,
    width,
    height,
    tamanho=14,
    cor=GRAFITE,
    espaco_linhas=1.15,
    espaco_antes=6,
    fonte=FONTE,
):
    """Insere varios paragrafos numa unica caixa.

    Cada item de ``linhas`` e uma string ou uma tupla ``(texto, opcoes)`` em que
    ``opcoes`` aceita as chaves ``tamanho``, ``cor``, ``negrito``, ``italico``,
    ``fonte``, ``recuo`` e ``marcador``.
    """
    caixa = slide.shapes.add_textbox(left, top, width, height)
    quadro = caixa.text_frame
    quadro.word_wrap = True
    quadro.margin_left = 0
    quadro.margin_right = 0
    quadro.margin_top = 0
    quadro.margin_bottom = 0

    primeiro = True
    for linha in linhas:
        if isinstance(linha, tuple):
            conteudo, opcoes = linha
        else:
            conteudo, opcoes = linha, {}

        paragrafo = quadro.paragraphs[0] if primeiro else quadro.add_paragraph()
        primeiro = False
        paragrafo.line_spacing = opcoes.get("espaco_linhas", espaco_linhas)
        paragrafo.space_before = Pt(0 if paragrafo is quadro.paragraphs[0] else espaco_antes)
        paragrafo.space_after = Pt(0)
        if opcoes.get("recuo"):
            paragrafo.level = opcoes["recuo"]

        marcador = opcoes.get("marcador", "")
        run = paragrafo.add_run()
        run.text = f"{marcador}{conteudo}" if marcador else conteudo
        run.font.size = Pt(opcoes.get("tamanho", tamanho))
        run.font.bold = opcoes.get("negrito", False)
        run.font.italic = opcoes.get("italico", False)
        run.font.color.rgb = opcoes.get("cor", cor)
        run.font.name = opcoes.get("fonte", fonte)
    return caixa


def cartao(slide, left, top, width, height, cor=BRANCO, borda=CINZA_CLARO, faixa=None):
    """Cartao retangular claro; ``faixa`` pinta uma barra vertical na borda esquerda."""
    shape = retangulo(slide, left, top, width, height, cor=cor, borda=borda)
    if faixa is not None:
        retangulo(slide, left, top, Inches(0.055), height, cor=faixa)
    return shape


ESPACO_LINHA_CODIGO = 1.06
PADDING_CODIGO = 0.26


def altura_codigo(n_linhas: int, tamanho: float = 11.5):
    """Altura necessaria para um bloco com ``n_linhas`` — evita transbordo."""
    altura_linha = tamanho * ESPACO_LINHA_CODIGO * 1.2 / 72
    return Inches(n_linhas * altura_linha + PADDING_CODIGO)


def codigo(slide, linhas_codigo, left, top, width, tamanho=11.5, destaque=None,
           height=None):
    """Bloco de codigo com fundo escuro e fonte monoespacada.

    A altura e calculada a partir do numero de linhas quando nao for informada,
    de modo que o texto nunca vaze do retangulo. Devolve a altura ocupada.

    ``destaque`` e um conjunto de indices de linha pintados em ambar (para
    apontar type hints, chamadas de NumPy, etc.).
    """
    destaque = destaque or set()
    if height is None:
        height = altura_codigo(len(linhas_codigo), tamanho)
    retangulo(slide, left, top, width, height, cor=RGBColor.from_string("14262F"))

    caixa = slide.shapes.add_textbox(
        left + Inches(0.16), top + Inches(0.12), width - Inches(0.32), height - Inches(0.2)
    )
    quadro = caixa.text_frame
    quadro.word_wrap = False
    quadro.margin_left = 0
    quadro.margin_right = 0
    quadro.margin_top = 0
    quadro.margin_bottom = 0

    for indice, linha in enumerate(linhas_codigo):
        paragrafo = quadro.paragraphs[0] if indice == 0 else quadro.add_paragraph()
        paragrafo.line_spacing = ESPACO_LINHA_CODIGO
        paragrafo.space_before = Pt(0)
        paragrafo.space_after = Pt(0)
        run = paragrafo.add_run()
        run.text = linha if linha else " "
        run.font.size = Pt(tamanho)
        run.font.name = FONTE_MONO
        if indice in destaque:
            run.font.color.rgb = RGBColor.from_string("F0B860")
            run.font.bold = True
        elif linha.strip().startswith("#") or linha.strip().startswith('"""'):
            run.font.color.rgb = RGBColor.from_string("7FA3B0")
        else:
            run.font.color.rgb = RGBColor.from_string("E4EBEF")
    return height


def tabela(
    slide,
    dados,
    left,
    top,
    width,
    height,
    larguras=None,
    tamanho_cabecalho=12,
    tamanho_corpo=11.5,
    cores_primeira_coluna=None,
    alinhamentos=None,
):
    """Tabela com cabecalho azul, zebra clara e bordas discretas.

    ``dados`` inclui a linha de cabecalho. ``larguras`` sao fracoes que somam 1.
    ``cores_primeira_coluna`` mapeia indice da linha de corpo -> cor do texto da
    primeira celula (usado para marcar status).
    """
    n_linhas = len(dados)
    n_colunas = len(dados[0])
    frame = slide.shapes.add_table(n_linhas, n_colunas, left, top, width, height)
    tabela_obj = frame.table

    if larguras:
        for indice, fracao in enumerate(larguras):
            tabela_obj.columns[indice].width = Emu(int(width * fracao))

    altura_cabecalho = Inches(0.36)
    tabela_obj.rows[0].height = altura_cabecalho
    altura_corpo = Emu(int((height - altura_cabecalho) / max(1, n_linhas - 1)))
    for linha in list(tabela_obj.rows)[1:]:
        linha.height = altura_corpo

    for i, linha_dados in enumerate(dados):
        for j, valor in enumerate(linha_dados):
            celula = tabela_obj.cell(i, j)
            celula.text = str(valor)
            celula.margin_left = Inches(0.09)
            celula.margin_right = Inches(0.07)
            celula.margin_top = Inches(0.03)
            celula.margin_bottom = Inches(0.03)
            celula.vertical_anchor = MSO_ANCHOR.MIDDLE

            celula.fill.solid()
            if i == 0:
                celula.fill.fore_color.rgb = AZUL
            else:
                celula.fill.fore_color.rgb = BRANCO if i % 2 else GELO_AZUL

            paragrafo = celula.text_frame.paragraphs[0]
            paragrafo.line_spacing = 0.95
            if alinhamentos and j < len(alinhamentos):
                paragrafo.alignment = alinhamentos[j]
            for run in paragrafo.runs:
                run.font.name = FONTE
                run.font.size = Pt(tamanho_cabecalho if i == 0 else tamanho_corpo)
                run.font.bold = i == 0 or j == 0
                if i == 0:
                    run.font.color.rgb = BRANCO
                elif j == 0 and cores_primeira_coluna and (i - 1) in cores_primeira_coluna:
                    run.font.color.rgb = cores_primeira_coluna[i - 1]
                else:
                    run.font.color.rgb = GRAFITE
    return tabela_obj


# ---------------------------------------------------------------------------
# Moldes de slide
# ---------------------------------------------------------------------------
def slide_em_branco(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_padrao(prs, titulo: str, kicker: str = "", numero: int | None = None):
    """Slide de conteudo: faixa fina no topo, kicker, titulo, regua ambar e rodape."""
    slide = slide_em_branco(prs)
    fundo(slide, GELO)

    retangulo(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.10), cor=AZUL)

    if kicker:
        texto(
            slide,
            kicker.upper(),
            MARGEM,
            Inches(0.42),
            LARGURA_UTIL,
            Inches(0.24),
            tamanho=10.5,
            cor=AZUL_CLARO,
            negrito=True,
        )

    texto(
        slide,
        titulo,
        MARGEM,
        Inches(0.68),
        LARGURA_UTIL,
        Inches(0.55),
        tamanho=27,
        cor=AZUL,
        negrito=True,
    )

    retangulo(slide, MARGEM, Inches(1.31), Inches(1.15), Inches(0.045), cor=AMBAR)

    if numero is not None:
        texto(
            slide,
            "Detector de Anomalias em Métricas de Servidores (SMD)  |  Grupo 12",
            MARGEM,
            Inches(7.03),
            Inches(9.0),
            Inches(0.26),
            tamanho=9,
            cor=CINZA,
        )
        texto(
            slide,
            f"{numero}/13",
            SLIDE_W - MARGEM - Inches(1.2),
            Inches(7.03),
            Inches(1.2),
            Inches(0.26),
            tamanho=9,
            cor=CINZA,
            alinhamento=PP_ALIGN.RIGHT,
        )
    return slide


def notas(slide, conteudo: str) -> None:
    """Grava as anotacoes do apresentador (painel de notas do PowerPoint)."""
    quadro = slide.notes_slide.notes_text_frame
    quadro.text = conteudo.strip()
    for paragrafo in quadro.paragraphs:
        for run in paragrafo.runs:
            run.font.size = Pt(12)
            run.font.name = FONTE
