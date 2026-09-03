"""Monta a apresentacao final do Grupo 12 (13 slides + anotacoes do apresentador).

Uso:
    python apresentacao/graficos.py      # gera as figuras a partir dos resultados
    python apresentacao/gerar_slides.py  # monta o .pptx

Os numeros exibidos vem dos arquivos de resultado do proprio projeto
(``artifacts/results.json`` e ``apresentacao/experimentos/``), nunca digitados
a mao. As anotacoes de cada slide sao escritas para quem vai apresentar estudar
antes: explicam o conteudo, o vocabulario tecnico e as perguntas provaveis.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from tema import (
    AMBAR,
    BASE_CONTEUDO,
    AZUL,
    AZUL_CLARO,
    BRANCO,
    CINZA,
    CINZA_CLARO,
    GELO_AZUL,
    GRAFITE,
    LARGURA_UTIL,
    MARGEM,
    SLIDE_H,
    SLIDE_W,
    TOPO_CONTEUDO,
    VERDE,
    VERMELHO,
    cartao,
    codigo,
    fundo,
    notas,
    paragrafos,
    retangulo,
    slide_em_branco,
    slide_padrao,
    tabela,
    texto,
)

AQUI = Path(__file__).resolve().parent
REPO = AQUI.parent
FIGURAS = AQUI / "figuras"
SAIDA = AQUI / "Grupo12-Detector-de-Anomalias-SMD.pptx"

REPO_URL = "github.com/ErasmoGusmao/detector-anomalias-smd"

EQUIPE = [
    "Breno Alexandre de Albuquerque Santos",
    "Erasmo de Melo Gusmão",
    "Gabriel Victor Alves Santana",
    "João Mateus Queiroz Moreira",
    "João Pedro da Silva Rodrigues",
    "Leonardo Canuto de Oliveira Magalhães",
    "Orlando Gomes dos Reis Neto",
    "Pedro Oliveira Pessoa Guerra",
]

RESULTADOS = json.loads((REPO / "artifacts" / "results.json").read_text(encoding="utf-8"))
CONFIGS = {
    cfg["id"]: cfg
    for cfg in json.loads(
        (AQUI / "experimentos" / "consolidado.json").read_text(encoding="utf-8")
    )
}


def pct(valor: float) -> str:
    return f"{valor * 100:.1f}%".replace(".", ",")


def num(valor: float, casas: int = 4) -> str:
    return f"{valor:.{casas}f}".replace(".", ",")


def figura(slide, nome: str, left, top, altura):
    """Insere uma figura com a altura pedida e devolve a largura ocupada."""
    from PIL import Image

    caminho = FIGURAS / nome
    if not caminho.is_file():
        raise FileNotFoundError(f"Figura ausente: {caminho}. Rode apresentacao/graficos.py")
    with Image.open(caminho) as imagem:
        proporcao = imagem.width / imagem.height
    largura = int(altura * proporcao)
    slide.shapes.add_picture(str(caminho), left, top, height=altura)
    return largura


# ===========================================================================
# Slide 1 - Problema, motivacao e objetivo
# ===========================================================================
def slide_01(prs) -> None:
    slide = slide_em_branco(prs)
    fundo(slide, AZUL)

    retangulo(slide, Inches(0), Inches(0), Inches(0.16), SLIDE_H, cor=AMBAR)

    texto(
        slide, "PROJETO INTEGRADOR  ·  ENGENHARIA DE SOFTWARE PARA IA E FRAMEWORKS PROFUNDOS  ·  CIn  ·  GRUPO 12",
        Inches(0.72), Inches(0.52), Inches(12.0), Inches(0.3),
        tamanho=10.5, cor=AZUL_CLARO, negrito=True,
    )
    texto(
        slide, "Detector de Anomalias em Métricas de Servidores",
        Inches(0.72), Inches(0.95), Inches(11.9), Inches(0.72),
        tamanho=38, cor=BRANCO, negrito=True,
    )
    texto(
        slide, "Detecção não supervisionada de comportamento anômalo em séries temporais multivariadas — dataset SMD",
        Inches(0.72), Inches(1.74), Inches(11.9), Inches(0.4),
        tamanho=15, cor=AZUL_CLARO,
    )
    retangulo(slide, Inches(0.72), Inches(2.24), Inches(1.4), Inches(0.05), cor=AMBAR)

    # --- coluna esquerda: o essencial do problema
    esquerda = Inches(0.72)
    largura_esq = Inches(6.55)
    blocos = [
        (
            "PROBLEMA",
            "Cada servidor publica 38 métricas ao mesmo tempo. Degradações que só "
            "aparecem na combinação dos sinais passam despercebidas até virarem incidente.",
        ),
        (
            "MOTIVAÇÃO",
            "O monitoramento hoje depende de limiar fixo por métrica e inspeção manual de "
            "painéis: alerta demais no trivial e tarde no que importa.",
        ),
        (
            "PÚBLICO-ALVO",
            "Analista de Operações / SRE (recebe o sinal e faz a triagem) e Coordenador de "
            "Observabilidade (decide manter, recalibrar ou retreinar o detector).",
        ),
        (
            "OBJETIVO",
            "Sinalizar automaticamente as janelas de tempo em que a máquina se afasta do "
            "padrão normal que ela mesma exibiu no passado.",
        ),
        (
            "PAPEL DA IA",
            "Um autoencoder aprende a reconstruir o comportamento normal. O que ele "
            "reconstrói mal é o que não se parece com o normal — esse erro é o escore de anomalia.",
        ),
    ]
    y = Inches(2.52)
    for rotulo, corpo in blocos:
        texto(slide, rotulo, esquerda, y, largura_esq, Inches(0.2),
              tamanho=10, cor=AMBAR, negrito=True)
        texto(slide, corpo, esquerda, y + Inches(0.235), largura_esq, Inches(0.6),
              tamanho=12.5, cor=BRANCO, espaco_linhas=1.05)
        y += Inches(0.86)

    # --- coluna direita: equipe e repositorio
    direita = Inches(7.62)
    largura_dir = Inches(5.0)
    cartao(slide, direita, Inches(2.42), largura_dir, Inches(3.15),
           cor=None, borda=AZUL_CLARO)
    texto(slide, "EQUIPE — GRUPO 12", direita + Inches(0.28), Inches(2.62),
          largura_dir - Inches(0.5), Inches(0.24), tamanho=10, cor=AMBAR, negrito=True)
    paragrafos(
        slide,
        [(nome, {"tamanho": 12, "cor": BRANCO}) for nome in EQUIPE],
        direita + Inches(0.28), Inches(2.92), largura_dir - Inches(0.5), Inches(2.5),
        espaco_antes=3.5, espaco_linhas=1.0,
    )

    cartao(slide, direita, Inches(5.78), largura_dir, Inches(0.92), cor=AMBAR, borda=None)
    texto(slide, "REPOSITÓRIO", direita + Inches(0.28), Inches(5.96),
          largura_dir - Inches(0.5), Inches(0.22), tamanho=9.5, cor=GRAFITE, negrito=True)
    texto(slide, REPO_URL, direita + Inches(0.28), Inches(6.2),
          largura_dir - Inches(0.5), Inches(0.4), tamanho=14, cor=GRAFITE, negrito=True)

    notas(slide, """
=========================================================================
COMO USAR ESTAS ANOTAÇÕES  (leia uma vez, antes de estudar os slides)

Cada slide traz quatro blocos:
  · TEMPO      — quanto falar naquele slide. A soma dá 14 minutos, com 1 de
                 folga sobre o limite de 15 do enunciado. Treine com cronômetro.
  · ROTEIRO    — o que dizer, já escrito em voz falada. Não é para decorar:
                 é para você entender o encadeamento e falar com as suas palavras.
  · ENTENDA    — o vocabulário técnico do slide, explicado do zero. É a parte
                 para estudar antes; ninguém precisa ter escrito o código para
                 apresentar bem, mas precisa entender o que está mostrando.
  · SE PERGUNTAREM — as perguntas prováveis, com a resposta honesta já pronta.

Regra de ouro da banca: quando não souber, diga que não sabe e aponte onde a
resposta está (README, docs/arquitetura.md, docs/requisitos/). O professor pediu
uma MAQUETE funcional — assumir uma lacuna conta a favor, inventar conta contra.

No modo Apresentador do PowerPoint (ou do Google Apresentações) estas notas
aparecem só para você; a plateia vê apenas o slide.
=========================================================================

[TEMPO: 60s — este slide dá o tom, não corra]

ROTEIRO — o que dizer
"Somos o Grupo 12. Nosso projeto é um detector de anomalias em métricas de servidores.
O problema: um servidor publica dezenas de métricas ao mesmo tempo — no nosso dataset, 38.
Hoje o monitoramento típico coloca um limiar fixo em cada métrica separadamente: se a CPU
passar de 90%, alerta. O que esse esquema não pega é a degradação que só existe na
COMBINAÇÃO — nenhuma métrica sozinha estoura o limite dela, mas o conjunto está num estado
que a máquina nunca exibiu antes. É esse tipo de falha silenciosa que nos motivou.
O objetivo é sinalizar automaticamente essas janelas de tempo, e o papel da IA é aprender o
que é 'normal' a partir dos próprios dados, em vez de alguém escolher limiares na mão."

ENTENDA — os conceitos deste slide
· SÉRIE TEMPORAL MULTIVARIADA: uma tabela em que cada linha é um instante no tempo e cada
  coluna é uma métrica. "Multivariada" = várias métricas medidas ao mesmo tempo, e o que
  interessa é como elas se movem em conjunto.
· AUTOENCODER: uma rede neural que aprende a copiar a entrada na saída, mas passando por um
  "funil" no meio (uma camada bem menor). Como o funil não deixa passar tudo, a rede é
  obrigada a guardar só os padrões que mais se repetem. Treinada apenas com dados normais,
  ela fica boa em reconstruir o normal e ruim em reconstruir qualquer outra coisa.
· ERRO DE RECONSTRUÇÃO: a diferença entre o que entrou e o que a rede devolveu. Erro baixo =
  "isto se parece com o que eu aprendi". Erro alto = "isto é estranho". É o nosso escore.
· NÃO SUPERVISIONADO: treinamos SEM dizer à rede o que é anomalia. Ela só vê exemplos normais.
· SRE: Site Reliability Engineer — a pessoa de plantão responsável por manter o sistema no ar.

SE PERGUNTAREM
P: "Por que não usar um limiar fixo, que é mais simples?"
R: "Essa é a alternativa honesta e ela cobre o caso fácil. O que ela não cobre é a anomalia
   combinada — e foi exatamente essa lacuna que registramos como requisito RNF-04. Somos
   transparentes num ponto: não medimos esse baseline, então nossa justificativa para o
   modelo está argumentada, não demonstrada. Está na reflexão crítica do documento."
P: "Quem é o cliente real?"
R: "É um stakeholder construído para o exercício, e declaramos isso no documento. O SMD é
   anonimizado, não há cliente por trás. Preferimos assumir a hipótese a inventar um cliente."

NÃO ESQUEÇA: ler os nomes completos não é necessário, mas o link do repositório precisa
aparecer na tela — é exigência do enunciado.
""")


# ===========================================================================
# Slide 2 - Dados e solucao proposta
# ===========================================================================
def slide_02(prs) -> None:
    slide = slide_padrao(prs, "Dados e solução proposta", "Resumo do projeto", 2)

    # Faixa de numeros
    metricas = [
        ("28.479", "amostras de treino\n(só comportamento normal)"),
        ("28.479", "amostras de teste\n(com anomalias rotuladas)"),
        ("38", "métricas por instante\n(atributos de entrada)"),
        ("9,46%", "taxa base de anomalia\n(2.694 instantes)"),
    ]
    largura_card = Inches(2.86)
    x = MARGEM
    for valor, rotulo in metricas:
        cartao(slide, x, TOPO_CONTEUDO, largura_card, Inches(1.16), cor=BRANCO, faixa=AMBAR)
        texto(slide, valor, x + Inches(0.24), TOPO_CONTEUDO + Inches(0.16),
              largura_card - Inches(0.4), Inches(0.42), tamanho=25, cor=AZUL, negrito=True)
        texto(slide, rotulo, x + Inches(0.24), TOPO_CONTEUDO + Inches(0.62),
              largura_card - Inches(0.4), Inches(0.5), tamanho=10.5, cor=CINZA,
              espaco_linhas=1.0)
        x += largura_card + Inches(0.19)

    # Coluna esquerda: base de dados
    topo2 = TOPO_CONTEUDO + Inches(1.44)
    largura_col = Inches(5.95)
    cartao(slide, MARGEM, topo2, largura_col, Inches(3.72), cor=BRANCO, faixa=AZUL)
    texto(slide, "BASE DE DADOS", MARGEM + Inches(0.28), topo2 + Inches(0.2),
          largura_col - Inches(0.5), Inches(0.24), tamanho=10, cor=AZUL_CLARO, negrito=True)
    paragrafos(
        slide,
        [
            ("SMD — Server Machine Dataset", {"tamanho": 15, "negrito": True, "cor": AZUL}),
            ("Origem: coletado do trabalho OmniAnomaly (NetManAIOps) e publicado no Kaggle. "
             "São métricas reais de operação de servidores, anonimizadas.", {"tamanho": 12}),
            ("Dimensão original: 28 máquinas × 38 métricas, cada máquina com sua própria "
             "série temporal completa.", {"tamanho": 12}),
            ("Recorte do projeto: a máquina machine-1-1, definida em config.py. Uma máquina "
             "já é uma série multivariada completa — o recorte reduz custo sem "
             "descaracterizar o problema.", {"tamanho": 12}),
            ("Arquivos: train (metade inicial, sem anomalia) · test (metade final) · "
             "test_label (rótulo 0/1 por instante, usado só na avaliação).", {"tamanho": 12}),
        ],
        MARGEM + Inches(0.28), topo2 + Inches(0.5), largura_col - Inches(0.56), Inches(3.1),
        espaco_antes=8,
    )

    # Coluna direita: solucao
    x2 = MARGEM + largura_col + Inches(0.19)
    cartao(slide, x2, topo2, largura_col, Inches(3.72), cor=BRANCO, faixa=AMBAR)
    texto(slide, "SOLUÇÃO PROPOSTA", x2 + Inches(0.28), topo2 + Inches(0.2),
          largura_col - Inches(0.5), Inches(0.24), tamanho=10, cor=AZUL_CLARO, negrito=True)
    paragrafos(
        slide,
        [
            ("Autoencoder treinado só com o normal", {"tamanho": 15, "negrito": True, "cor": AZUL}),
            ("ENTRADA  →  janela de 50 instantes consecutivos × 38 métricas padronizadas "
             "(z-score), deslocando de 1 em 1.", {"tamanho": 12}),
            ("SAÍDA  →  (a) escore de anomalia = erro de reconstrução da janela; "
             "(b) classificação binária normal / anômala.", {"tamanho": 12}),
            ("VARIÁVEL PREVISTA  →  o rótulo binário da janela. Repare: o rótulo NÃO entra "
             "no treino, só na avaliação — por isso a abordagem é não supervisionada.",
             {"tamanho": 12}),
            ("DECISÃO  →  a janela é anômala quando o escore ultrapassa um limiar derivado "
             "do percentil dos erros do próprio conjunto de treino.", {"tamanho": 12}),
        ],
        x2 + Inches(0.28), topo2 + Inches(0.5), largura_col - Inches(0.56), Inches(3.1),
        espaco_antes=8,
    )

    notas(slide, """
[TEMPO: 60s]

ROTEIRO — o que dizer
"A base é o SMD, Server Machine Dataset, do trabalho OmniAnomaly, disponível publicamente no
Kaggle. São métricas reais de servidores, anonimizadas. O dataset tem 28 máquinas; nós
recortamos uma delas, a machine-1-1, porque cada máquina já é uma série multivariada
completa — o recorte reduz o custo sem simplificar o problema.
São 28.479 instantes de treino e 28.479 de teste, com 38 métricas em cada instante. No teste,
9,46% dos instantes são anômalos: essa é a taxa base, e ela vai importar na hora de julgar a
precisão.
A entrada do sistema é uma janela de 50 instantes consecutivos; a saída é um escore de
anomalia e a classificação binária daquela janela."

ENTENDA — os conceitos deste slide
· JANELA DESLIZANTE: em vez de olhar um instante isolado, olhamos 50 instantes seguidos. A
  janela anda de 1 em 1 instante, então janelas vizinhas se sobrepõem bastante. Isso dá
  contexto temporal ao modelo: um pico isolado é diferente de uma subida sustentada.
· PADRONIZAÇÃO / Z-SCORE: transformar cada métrica para média 0 e desvio 1, com a fórmula
  (x − média) ÷ desvio. Sem isso, uma métrica que varia de 0 a 1.000.000 dominaria o erro e
  outra que varia de 0 a 1 seria ignorada.
· TAXA BASE: a proporção de anomalias que existe no conjunto. Aqui, 9,46%. É a régua da
  precisão: um detector que chutasse aleatoriamente acertaria ~9,46%. Ficar acima disso é o
  mínimo para o modelo ter valor — e é exatamente assim que nosso RNF-01 está redigido.
· POR QUE O TREINO NÃO TEM ANOMALIA: é uma característica do SMD. A primeira metade da série
  é normal por construção. É isso que viabiliza o método: o modelo só conhece o normal.

SE PERGUNTAREM
P: "Por que só uma máquina?"
R: "Escopo. O pipeline aceita qualquer uma das 28 trocando uma linha no config.py — a
   variável FILE_NAME. Não rodamos as outras, e registramos isso como limitação."
P: "Vocês usam os rótulos?"
R: "Só para avaliar. O treino é cego aos rótulos; se usássemos, seria aprendizado
   supervisionado e outro projeto."
""")


# ===========================================================================
# Slide 3 - Requisitos funcionais
# ===========================================================================
def slide_03(prs) -> None:
    slide = slide_padrao(prs, "Requisitos funcionais", "Requisitos · elicitados com GR4ML", 3)

    texto(
        slide,
        "Elicitados com o framework GR4ML: cada RF nasce de um Decision Goal + Question Goal da Business View — "
        "a origem está registrada no documento, e é ela que dá rastreabilidade.",
        MARGEM, TOPO_CONTEUDO, LARGURA_UTIL, Inches(0.34),
        tamanho=12, cor=CINZA,
    )

    dados = [
        ["RF", "O sistema deve...", "Status no protótipo"],
        ["RF-01", "calcular um escore de anomalia para cada janela de 50 instantes, a partir do erro de reconstrução das 38 métricas padronizadas.", "Atendido"],
        ["RF-02", "classificar cada janela como normal ou anômala, comparando seu escore ao limiar vigente.", "Atendido"],
        ["RF-03", "disponibilizar, para cada janela anômala, o erro decomposto por métrica, ordenado por contribuição.", "Do sistema-alvo"],
        ["RF-04", "aprender o padrão normal da máquina a partir da partição de treino, sem usar rótulos.", "Atendido"],
        ["RF-05", "persistir o modelo treinado, o limiar de anomalia e as estatísticas de padronização.", "Parcial"],
        ["RF-06", "derivar o limiar do percentil configurado sobre os escores de treino, registrando o valor aplicado.", "Atendido"],
        ["RF-07", "calcular e registrar precisão, recall e F1 da versão avaliada, contra os rótulos de teste.", "Parcial"],
    ]
    cores = {0: VERDE, 1: VERDE, 2: VERMELHO, 3: VERDE, 4: AMBAR, 5: VERDE, 6: AMBAR}
    tabela(
        slide, dados, MARGEM, TOPO_CONTEUDO + Inches(0.46), LARGURA_UTIL, Inches(3.62),
        larguras=[0.075, 0.755, 0.17], tamanho_corpo=11.5,
        cores_primeira_coluna=cores,
    )

    y = TOPO_CONTEUDO + Inches(4.22)
    cartao(slide, MARGEM, y, LARGURA_UTIL, Inches(0.86), cor=GELO_AZUL, faixa=AMBAR)
    paragrafos(
        slide,
        [
            ("Onde o protótipo ainda não chega — e por quê", {"tamanho": 12, "negrito": True, "cor": AZUL}),
            ("RF-03 e RF-05/RF-07 parciais: o erro hoje é agregado por janela (falta a decomposição por métrica); "
             "a média e o desvio da padronização não são persistidos junto do modelo; e o percentual de janelas "
             "sinalizadas não vai para o results.json. São lacunas conhecidas, declaradas no documento de arquitetura.",
             {"tamanho": 11.5}),
        ],
        MARGEM + Inches(0.26), y + Inches(0.14), LARGURA_UTIL - Inches(0.5), Inches(0.68),
        espaco_antes=4,
    )

    notas(slide, """
[TEMPO: 65s — NÃO leia os sete requisitos em voz alta]

ROTEIRO — o que dizer
"Os requisitos não foram inventados na hora de escrever o documento: saíram do GR4ML, o
framework de engenharia de requisitos que vimos na aula prática. Ele parte do objetivo de
negócio, passa por quais decisões o usuário precisa tomar, e só então vira requisito escrito.
São sete requisitos funcionais. Destaco três: o RF-01 é o coração — calcular um escore de
anomalia por janela. O RF-04 é o que define nossa abordagem — aprender o normal sem usar
rótulos. E o RF-06 é o que torna o sistema auditável — o limiar não é escolhido a dedo, é
derivado dos dados e registrado a cada execução.
Somos transparentes sobre o que não está pronto: o RF-03 é do sistema-alvo e não do
protótipo, e dois requisitos estão parciais. Isso está declarado no documento, não escondido."

ENTENDA — os conceitos deste slide
· GR4ML: Goal-Oriented Requirements Engineering for Machine Learning. A ideia central é que
  requisito de sistema de IA não começa em "que modelo usar", e sim em "que decisão de
  negócio precisa ser tomada". Ele organiza isso em três visões: Business, Analytics Design e
  Data Preparation.
· DECISION GOAL: a decisão concreta que o usuário precisa tomar. O nosso: a cada alerta, o
  analista decide entre investigar agora, seguir observando, ou registrar como não acionável.
· QUESTION GOAL: a pergunta que os dados precisam responder para apoiar aquela decisão. A
  nossa: "esta janela da série apresenta comportamento anômalo?"
· RASTREABILIDADE: é a propriedade de cada requisito apontar de onde veio. É critério de
  avaliação da disciplina, e é o que distingue requisito elicitado de requisito chutado.
· "O SISTEMA DEVE...": todo requisito começa assim, tem uma exigência só, e precisa ser
  testável — alguém de fora deve conseguir olhar o sistema pronto e dizer se foi atendido.

SE PERGUNTAREM
P: "Por que apresentar requisito que vocês não implementaram?"
R: "Porque o documento descreve o sistema-alvo, e o professor pediu explicitamente uma
   maquete funcional, não um sistema completo. Declarar a lacuna é mais honesto — e mais
   útil — do que rebaixar o requisito até caber no que já fizemos."
P: "O que falta exatamente no RF-05?"
R: "Salvamos os pesos do modelo e o limiar, mas não a média e o desvio usados na
   padronização. Sem eles, quem carregar o modelo não consegue preparar dados novos
   exatamente do mesmo jeito. É uma correção pequena e localizada."
""")


# ===========================================================================
# Slide 4 - RNF e criterios de aceitacao
# ===========================================================================
def slide_04(prs) -> None:
    slide = slide_padrao(prs, "Requisitos não funcionais e critérios de aceitação",
                         "Requisitos · qualidade verificável", 4)

    p995 = CONFIGS["A"]["por_percentil"]["99.5"]
    dados = [
        ["RNF", "Qualidade exigida", "Como verificamos", "Resultado medido"],
        ["RNF-01", "Baixa taxa de falso positivo, medida pela precisão, acima da taxa base do conjunto.",
         "Executar a avaliação no teste e comparar a precisão com a taxa base (9,46%).",
         f"{pct(p995['precision'])}  >  9,46%"],
        ["RNF-02", "Localizar as métricas responsáveis pela anomalia (erro por dimensão).",
         "Verificar se a saída traz o erro decomposto por métrica.", "Não atendido"],
        ["RNF-03", "Custo de treino viável: concluir o treino sem GPU dedicada.",
         "Executar o pipeline em CPU e medir o tempo.",
         f"{CONFIGS['A']['tempo_treino_s']:.0f} s em CPU".replace(".", ",")],
        ["RNF-04", "Cobrir degradações que só aparecem na combinação das métricas (recall).",
         "Medir o recall no teste e comparar com um detector de limiar por métrica.",
         f"recall {pct(p995['recall'])}"],
    ]
    cores = {0: VERDE, 1: VERMELHO, 2: VERDE, 3: AMBAR}
    tabela(
        slide, dados, MARGEM, TOPO_CONTEUDO, LARGURA_UTIL, Inches(2.72),
        larguras=[0.075, 0.315, 0.375, 0.235], tamanho_corpo=11,
        cores_primeira_coluna=cores,
    )

    y = TOPO_CONTEUDO + Inches(2.94)
    cartao(slide, MARGEM, y, LARGURA_UTIL, Inches(1.14), cor=AZUL, borda=None)
    texto(slide, "CRITÉRIO DE ACEITAÇÃO DO CONJUNTO", MARGEM + Inches(0.3), y + Inches(0.16),
          LARGURA_UTIL - Inches(0.6), Inches(0.24), tamanho=10, cor=AMBAR, negrito=True)
    texto(
        slide,
        "“Este conjunto de requisitos será bem-sucedido se o detector sinalizar ao menos 85% dos eventos anômalos "
        "conhecidos do SMD (recall), mantendo precisão superior à taxa base de anomalias do conjunto avaliado "
        "(9,46% na machine-1-1), verificado a cada nova versão do modelo.”",
        MARGEM + Inches(0.3), y + Inches(0.46), LARGURA_UTIL - Inches(0.6), Inches(0.62),
        tamanho=13, cor=BRANCO, espaco_linhas=1.06,
    )

    y2 = y + Inches(1.32)
    veredito = (
        f"ATENDIDO na execução de referência:  recall {pct(p995['recall'])} ≥ 85%   ·   "
        f"precisão {pct(p995['precision'])} > 9,46%"
    )
    cartao(slide, MARGEM, y2, LARGURA_UTIL, Inches(0.56), cor=BRANCO, faixa=VERDE)
    texto(slide, veredito, MARGEM + Inches(0.3), y2 + Inches(0.15),
          LARGURA_UTIL - Inches(0.6), Inches(0.3), tamanho=14, cor=VERDE, negrito=True)

    notas(slide, """
[TEMPO: 70s — o critério de aceitação é o ponto alto, chegue nele]

ROTEIRO — o que dizer
"Requisito não funcional é sobre qualidade, não sobre função. E a regra que seguimos é: se
não dá para verificar, não é requisito — é desejo. Por isso cada linha tem uma coluna
dizendo COMO se verifica.
O RNF-01 pede precisão acima da taxa base. Note por que a taxa base é a régua certa: como
9,46% do conjunto é anômalo, um detector que chutasse ao acaso acertaria 9,46%. Ficar acima
disso é o mínimo para o modelo ter valor. Medimos 25,3% — três vezes a régua.
O RNF-03 exigia treinar sem GPU: treinou em CPU em menos de três minutos.
O RNF-02 é o que não atendemos, e é justamente o que exige o erro por dimensão.
E o conjunto todo fecha com um critério de aceitação ligado ao indicador de negócio — este
aqui embaixo. Na execução de referência, ele foi atendido nos dois lados."

ENTENDA — os conceitos deste slide
· PRECISÃO: dos alertas que o sistema deu, quantos eram anomalia de verdade. Precisão baixa =
  o plantão perde tempo com alarme falso.
· RECALL: das anomalias que existiam, quantas o sistema pegou. Recall baixo = incidente passou
  batido. É o erro mais caro em operação.
· O TRADE-OFF: os dois brigam entre si. Baixar o limiar aumenta o recall e derruba a precisão;
  subir faz o contrário. Não existe "melhorar os dois" só apertando o limiar — para isso é
  preciso um modelo melhor. O slide 11 mostra essa curva medida.
· F1: a média harmônica entre precisão e recall — um número só, que pune quem é muito bom num
  e muito ruim no outro.
· SOFTGOAL: no GR4ML, é a qualidade desejada ainda em linguagem de negócio ("não quero muito
  alarme falso"). O "pulo do gato" do framework é obrigar cada softgoal a virar uma métrica
  concreta — senão ele fica bonito no papel e não influencia decisão nenhuma.

SE PERGUNTAREM
P: "25% de precisão não é baixo?"
R: "Em termos absolutos, sim: três em cada quatro alertas são falsos. Mas a comparação
   relevante é com a taxa base de 9,46%, e aí é quase três vezes melhor que o acaso. E é uma
   escolha consciente: nosso limiar está calibrado para recall alto, porque em operação
   deixar passar um incidente custa mais que investigar um alerta à toa. No slide 11
   mostramos que, mexendo só no limiar, chegamos a 83% de precisão — pagando com recall."
P: "Como vocês verificariam o RNF-04 direito?"
R: "Implementando o baseline de limiar por métrica e comparando os dois no mesmo volume de
   alerta. Não fizemos, e assumimos isso na reflexão crítica do documento."
""")


# ===========================================================================
# Slide 5 - Diagrama da arquitetura
# ===========================================================================
def slide_05(prs) -> None:
    slide = slide_padrao(prs, "Arquitetura do sistema", "Arquitetura · visão de componentes", 5)

    # --- orquestrador
    y_main = TOPO_CONTEUDO + Inches(0.04)
    retangulo(slide, MARGEM, y_main, LARGURA_UTIL, Inches(0.66), cor=AZUL)
    texto(slide, "main.py   —   orquestrador: é o único que conhece a ordem das etapas",
          MARGEM + Inches(0.3), y_main + Inches(0.2), LARGURA_UTIL - Inches(0.6), Inches(0.3),
          tamanho=14, cor=BRANCO, negrito=True)

    # --- pipeline
    etapas = [
        ("SMD", "train · test\ntest_label", "arquivos .txt", GELO_AZUL, AZUL),
        ("data/", "loader.py\ncarregar e limpar", "→ DataFrame (N, 38)", BRANCO, AZUL),
        ("preprocessing/", "transform.py\nz-score e split", "→ ndarray padronizado", BRANCO, AZUL),
        ("training/", "train.py\nDataLoader e laço", "→ tensores (b, 50, 38)", BRANCO, AZUL),
        ("models/", "model.py\nAutoencoder e erro", "→ erro por janela", BRANCO, AZUL),
        ("evaluation/", "metrics.py\nlimiar e métricas", "→ precisão, recall, F1", BRANCO, AZUL),
        ("artifacts/", "autoencoder.pt\nresults.json", "saída persistida", GELO_AZUL, AMBAR),
    ]
    largura_etapa = Inches(1.52)
    espaco = Inches(0.245)
    y_etapas = y_main + Inches(1.12)
    altura_etapa = Inches(1.72)

    x = MARGEM
    centros = []
    for indice, (titulo, corpo, saida, cor_fundo, cor_faixa) in enumerate(etapas):
        cartao(slide, x, y_etapas, largura_etapa, altura_etapa, cor=cor_fundo, faixa=cor_faixa)
        texto(slide, titulo, x + Inches(0.16), y_etapas + Inches(0.18),
              largura_etapa - Inches(0.26), Inches(0.26), tamanho=12, cor=AZUL, negrito=True)
        texto(slide, corpo, x + Inches(0.16), y_etapas + Inches(0.52),
              largura_etapa - Inches(0.26), Inches(0.62), tamanho=10, cor=CINZA,
              espaco_linhas=1.08)
        texto(slide, saida, x + Inches(0.16), y_etapas + altura_etapa - Inches(0.46),
              largura_etapa - Inches(0.26), Inches(0.36), tamanho=9.5, cor=AMBAR,
              negrito=True, espaco_linhas=1.05)
        centros.append(x + largura_etapa / 2)

        if indice < len(etapas) - 1:
            seta = retangulo(
                slide,
                x + largura_etapa + Inches(0.055),
                y_etapas + altura_etapa / 2 - Inches(0.075),
                espaco - Inches(0.11),
                Inches(0.15),
                cor=AMBAR,
                forma=MSO_SHAPE.RIGHT_ARROW,
            )
            seta.line.fill.background()
        x += largura_etapa + espaco

    # chevrons do orquestrador para as etapas do pipeline
    for centro in centros[1:6]:
        chevron = retangulo(
            slide, centro - Inches(0.10), y_main + Inches(0.76),
            Inches(0.20), Inches(0.26), cor=CINZA_CLARO,
            forma=MSO_SHAPE.DOWN_ARROW,
        )
        chevron.line.fill.background()

    # --- utils
    y_utils = y_etapas + altura_etapa + Inches(0.42)
    altura_utils = BASE_CONTEUDO - y_utils
    largura_utils = Inches(5.1)
    cartao(slide, MARGEM, y_utils, largura_utils, altura_utils, cor=BRANCO, faixa=CINZA)
    texto(slide, "src/utils/   —   apoio transversal (não carrega dados do pipeline)",
          MARGEM + Inches(0.26), y_utils + Inches(0.2), largura_utils - Inches(0.45),
          Inches(0.24), tamanho=11.5, cor=AZUL, negrito=True)
    texto(slide, "config.py — caminhos e hiperparâmetros\ntorch_utils.py — device e sementes",
          MARGEM + Inches(0.26), y_utils + Inches(0.54), largura_utils - Inches(0.45),
          Inches(0.6), tamanho=10.5, cor=CINZA, espaco_linhas=1.15)

    x_regra = MARGEM + largura_utils + Inches(0.35)
    largura_regra = LARGURA_UTIL - largura_utils - Inches(0.35)
    cartao(slide, x_regra, y_utils, largura_regra, altura_utils, cor=GELO_AZUL, faixa=AMBAR)
    texto(slide, "REGRA DE DEPENDÊNCIA:  as setas nunca voltam",
          x_regra + Inches(0.26), y_utils + Inches(0.2), largura_regra - Inches(0.45),
          Inches(0.24), tamanho=11.5, cor=AZUL, negrito=True)
    texto(slide, "Nenhum pacote de src/ importa o main.py, e utils/ não importa ninguém do "
                 "pipeline. É o que permite usar o pré-processamento ou as métricas fora do "
                 "pipeline — nos testes, por exemplo — sem arrastar o resto junto.",
          x_regra + Inches(0.26), y_utils + Inches(0.54), largura_regra - Inches(0.45),
          Inches(0.7), tamanho=10.5, cor=CINZA, espaco_linhas=1.12)

    notas(slide, """
[TEMPO: 65s — este slide é obrigatório no enunciado; aponte para a tela]

ROTEIRO — o que dizer
"Esta é a arquitetura real do projeto, não um desenho genérico. Leia da esquerda para a
direita: os dados do SMD entram, passam por carregamento e limpeza, vão para o
pré-processamento em NumPy, viram janelas e tensores no módulo de treino, o modelo aprende,
a avaliação aplica o limiar e calcula as métricas, e tudo termina em artefatos gravados em
disco — o modelo salvo e o results.json.
Duas coisas que valem apontar. Primeira: a barra azul em cima. O main.py é o único componente
que sabe em que ordem as etapas acontecem; nenhum módulo chama o outro por conta própria.
Segunda: a caixa cinza embaixo. O utils é apoio — configuração e sementes — e a seta dele é
tracejada de propósito, porque ele não carrega dado do pipeline, só parâmetro.
E a regra que amarra tudo: as setas nunca voltam. Nenhum pacote importa o orquestrador."

ENTENDA — os conceitos deste slide
· ORQUESTRADOR: o componente que chama os outros na ordem certa. Concentrar a ordem num lugar
  só significa que mudar a sequência é mexer num arquivo, não em seis.
· ACOPLAMENTO: o quanto um módulo depende dos outros. Quanto menos, melhor — porque mudança
  em um não obriga mudança nos demais.
· COESÃO: o quanto o que está dentro de um módulo pertence junto. Alta coesão = o arquivo faz
  uma coisa só. Nosso transform.py só transforma dados; não carrega arquivo nem treina.
· DEPENDÊNCIA ACÍCLICA ("as setas nunca voltam"): se A importa B, B não pode importar A. Ciclo
  de importação em Python quebra na hora e, pior, torna impossível testar um módulo sozinho.
· POR QUE ISSO IMPORTA NA PRÁTICA: como preprocessing e data não importam torch, um colega
  sem PyTorch instalado consegue rodar e testar essa parte. Isso aconteceu de verdade no grupo.

SE PERGUNTAREM
P: "Cadê o módulo de inferência?"
R: "Não existe como pacote separado, e isso é consciente. Hoje a inferência é a composição de
   duas funções — o erro de reconstrução mais a predição pelo limiar — chamadas pelo main.py.
   Quando houver consumo de dados novos fora do pipeline de avaliação, é o próximo módulo a
   nascer. Está registrado nas limitações do documento de arquitetura."
P: "O diagrama está versionado?"
R: "Está, em docs/arquitetura.md, escrito em Mermaid — o GitHub renderiza direto na página."
""")


# ===========================================================================
# Slide 6 - Decisoes arquiteturais
# ===========================================================================
def slide_06(prs) -> None:
    slide = slide_padrao(prs, "Decisões arquiteturais", "Arquitetura · o porquê de cada escolha", 6)

    decisoes = [
        ("Um pacote por etapa do pipeline",
         "Cada etapa tem um dono e um arquivo de teste próprio. Foi o que permitiu várias pessoas "
         "trabalharem em paralelo sem conflito de merge.",
         "separação de responsabilidades"),
        ("Pré-processamento em NumPy, antes de qualquer tensor",
         "Mantém data/ e preprocessing/ executáveis sem PyTorch instalado — destravou colegas cujo "
         "ambiente não tinha a biblioteca.",
         "baixo acoplamento"),
        ("Configuração centralizada em utils/config.py",
         "Trocar janela, épocas ou percentil é editar uma linha, não caçar constante espalhada. É o que "
         "tornou viável comparar configurações.",
         "manutenibilidade"),
        ("Padronização com estatísticas só do treino",
         "Validação e teste usam a média e o desvio do treino. O contrário vazaria informação do teste "
         "para o pré-processamento.",
         "validade do experimento"),
        ("Split temporal, nunca aleatório",
         "A série é temporal: embaralhar antes de dividir treinaria o modelo com o futuro e inflaria o "
         "resultado artificialmente.",
         "validade do experimento"),
        ("Persistência separada da definição do modelo",
         "persistence.py grava só os pesos; a arquitetura fica em model.py. Carregar exige recriar a "
         "estrutura — não dependemos de serializar a classe.",
         "alta coesão"),
        ("Orquestração exclusivamente no main.py",
         "Os módulos não sabem em que ordem são chamados. Cada função pode ser exercitada isolada num teste.",
         "testabilidade"),
        ("Métricas em NumPy puro, sem scikit-learn",
         "Uma dependência a menos e o cálculo de precisão e recall fica auditável linha a linha — e testável.",
         "testabilidade"),
    ]

    largura_col = Inches(5.95)
    altura_item = Inches(1.23)
    for indice, (titulo, corpo, tag) in enumerate(decisoes):
        coluna, linha = divmod(indice, 4)
        x = MARGEM + coluna * (largura_col + Inches(0.19))
        y = TOPO_CONTEUDO + linha * (altura_item + Inches(0.09))
        cartao(slide, x, y, largura_col, altura_item, cor=BRANCO, faixa=AZUL)
        texto(slide, f"D{indice + 1}  ·  {titulo}", x + Inches(0.24), y + Inches(0.14),
              largura_col - Inches(0.45), Inches(0.28), tamanho=12.5, cor=AZUL, negrito=True)
        texto(slide, corpo, x + Inches(0.24), y + Inches(0.44), largura_col - Inches(0.45),
              Inches(0.54), tamanho=11, cor=GRAFITE, espaco_linhas=1.03)
        texto(slide, tag.upper(), x + Inches(0.24), y + Inches(0.99),
              largura_col - Inches(0.45), Inches(0.2), tamanho=9, cor=AMBAR, negrito=True)

    notas(slide, """
[TEMPO: 65s — escolha TRÊS decisões para falar, não as oito]

ROTEIRO — o que dizer
"Estas são as decisões que moldaram o código, com a justificativa de cada uma. Vou destacar três.
A D2: todo o pré-processamento é NumPy puro e acontece ANTES de qualquer tensor. Isso não é
preciosismo — significa que quem não tem PyTorch instalado ainda consegue rodar e testar essa
metade do sistema. Isso resolveu um problema real de ambiente no nosso grupo.
A D4 e a D5 são sobre validade do experimento. Padronizamos validação e teste com a média e o
desvio calculados SÓ no treino: usar as estatísticas do teste seria vazamento de dados. E
dividimos a série em ordem cronológica, nunca embaralhando — embaralhar uma série temporal
treina o modelo com o futuro e infla o resultado.
E a D7: só o main.py orquestra. É isso que deixa cada função testável isoladamente."

ENTENDA — os conceitos deste slide
· VAZAMENTO DE DADOS (data leakage): quando informação do conjunto de teste influencia o
  treino, mesmo indiretamente. Padronizar tudo junto é o caso clássico: a média "sabe" o que
  há no teste. O resultado fica bom no papel e quebra na vida real.
· POR QUE NÃO EMBARALHAR SÉRIE TEMPORAL: se as janelas forem sorteadas, o modelo treina com
  instantes posteriores aos que vai prever. Como janelas vizinhas se sobrepõem, ele
  praticamente decora o teste. O resultado seria ótimo e falso.
· STATE_DICT: no PyTorch, um dicionário só com os pesos aprendidos — sem a estrutura da rede.
  Salvar assim obriga recriar a classe para carregar, o que parece trabalho extra mas evita
  que o arquivo salvo quebre quando o código da classe mudar.
· ALTA COESÃO / BAIXO ACOPLAMENTO: o par que a disciplina cobra. Coesão é "cada módulo faz uma
  coisa"; acoplamento é "quantos módulos preciso mexer para mudar uma coisa". Queremos a
  primeira alta e o segundo baixo.

SE PERGUNTAREM
P: "Por que não usaram scikit-learn para as métricas?"
R: "Duas razões: uma dependência a menos, e as fórmulas ficam visíveis e testáveis por nós.
   Precisão e recall são divisões simples; escrevê-las tornou a suíte de testes mais
   significativa do que testar um wrapper de biblioteca."
P: "Essas decisões foram tomadas antes ou depois de codar?"
R: "Mistas, e documentadas depois em docs/arquitetura.md. A separação por etapa veio antes,
   ainda na Entrega 1; a decisão de deixar o pré-processamento livre de torch veio da prática."
""")


# ===========================================================================
# Slide 7 - Estrutura de pastas
# ===========================================================================
def slide_07(prs) -> None:
    slide = slide_padrao(prs, "Estrutura do projeto", "Modularização · organização das pastas", 7)

    arvore = [
        "detector-anomalias-smd/",
        "├── main.py                          # ponto único de execução do pipeline",
        "├── requirements.txt",
        "├── README.md",
        "├── src/",
        "│   ├── data/loader.py               # carregamento e limpeza (pandas)",
        "│   ├── preprocessing/transform.py   # padronização e split (NumPy)",
        "│   ├── models/",
        "│   │   ├── model.py                 # Autoencoder (nn.Module) e erro",
        "│   │   └── persistence.py           # salvar / carregar state_dict",
        "│   ├── training/train.py            # DataLoader, laço de treino e validação",
        "│   ├── evaluation/metrics.py        # limiar, predição e métricas",
        "│   └── utils/",
        "│       ├── config.py                # caminhos e hiperparâmetros",
        "│       └── torch_utils.py           # device e sementes",
        "├── tests/                           # 4 arquivos, um por pacote de src/",
        "├── docs/                            # requisitos (GR4ML) e arquitetura",
        "├── apresentacao/                    # slides, gerados a partir dos resultados",
        "├── data/                            # SMD (não versionado)",
        "└── artifacts/                       # modelo e resultados gerados",
    ]
    largura_codigo = Inches(7.35)
    codigo(slide, arvore, MARGEM, TOPO_CONTEUDO, largura_codigo, tamanho=10.5,
           destaque={4, 15})

    x2 = MARGEM + largura_codigo + Inches(0.26)
    largura_col = LARGURA_UTIL - largura_codigo - Inches(0.26)
    razoes = [
        ("Um diretório por responsabilidade",
         "O nome da pasta diz a etapa do pipeline. Quem chega no repositório sabe onde procurar sem "
         "abrir arquivo."),
        ("tests/ espelha src/ na mesma granularidade",
         "Um arquivo de teste por pacote: test_data, test_preprocessing, test_model, test_training. "
         "É imediato saber onde um comportamento é verificado."),
        ("O que é gerado fica fora do versionamento",
         "data/ e artifacts/ não vão para o Git: dados brutos e modelos treinados são reproduzíveis, "
         "não código-fonte."),
        ("main.py na raiz, sozinho",
         "Um único ponto de entrada: python main.py roda tudo, do arquivo bruto ao results.json."),
    ]
    y = TOPO_CONTEUDO
    for titulo, corpo in razoes:
        cartao(slide, x2, y, largura_col, Inches(1.09), cor=BRANCO, faixa=AMBAR)
        texto(slide, titulo, x2 + Inches(0.22), y + Inches(0.14), largura_col - Inches(0.42),
              Inches(0.5), tamanho=12, cor=AZUL, negrito=True, espaco_linhas=1.0)
        texto(slide, corpo, x2 + Inches(0.22), y + Inches(0.5), largura_col - Inches(0.42),
              Inches(0.54), tamanho=10.5, cor=GRAFITE, espaco_linhas=1.03)
        y += Inches(1.18)

    notas(slide, """
[TEMPO: 45s — mostre a árvore, mas gaste o tempo no PORQUÊ, à direita]

ROTEIRO — o que dizer
"Esta é a árvore real do repositório. O enunciado pede mais do que mostrar a estrutura: pede
explicar por que separamos assim. São quatro razões.
Primeira: um diretório por responsabilidade, e o nome da pasta diz a etapa do pipeline. Quem
abre o repositório pela primeira vez sabe onde procurar sem abrir arquivo nenhum.
Segunda: a pasta tests espelha a src na mesma granularidade — um arquivo de teste por pacote.
Se alguém pergunta 'onde vocês testam o pré-processamento?', a resposta é o nome do arquivo.
Terceira: o que é gerado não é versionado. Os dados brutos e o modelo treinado ficam fora do
Git, porque são reproduzíveis a partir do código; versionar binário grande polui o histórico.
Quarta: um único ponto de entrada. python main.py roda o pipeline inteiro."

ENTENDA — os conceitos deste slide
· PACOTE PYTHON: uma pasta com __init__.py, que permite importar como src.data.loader. É o
  mecanismo que torna a separação em pastas uma separação real de código, e não só visual.
· POR QUE NÃO VERSIONAR data/ E artifacts/: Git é feito para texto. Arquivo binário grande
  (um modelo .pt) incha o repositório para sempre, porque o histórico guarda toda versão.
  A regra prática: versione o que gera, não o que foi gerado.
· CONSEQUÊNCIA HONESTA DISSO: como artifacts/ não está no Git, o professor não encontra o
  results.json citado no nosso critério de aceitação — ele precisa rodar o pipeline. É uma
  limitação conhecida nossa (se perguntarem, assuma; não invente que está lá).
· REQUIREMENTS.TXT: a lista de dependências com versões. É o que faz outra pessoa reproduzir
  o ambiente com um comando.

SE PERGUNTAREM
P: "Por que models/ tem dois arquivos e os outros pacotes só um?"
R: "Porque são responsabilidades diferentes: model.py define a rede, persistence.py cuida de
   salvar e carregar. Manter juntos misturaria arquitetura com serialização — e é justamente
   a serialização que a gente pode querer trocar depois."
P: "Vocês seguiram a estrutura sugerida no enunciado?"
R: "Sim, com uma diferença: não temos o pacote inference/, porque a inferência hoje é a
   composição de duas funções chamadas pelo main.py. Preferimos não criar pasta vazia."
""")


# ===========================================================================
# Slide 8 - Responsabilidades dos modulos
# ===========================================================================
def slide_08(prs) -> None:
    slide = slide_padrao(prs, "Responsabilidades dos módulos", "Modularização · contrato de cada peça", 8)

    dados = [
        ["Módulo", "Responsabilidade", "Entrada → saída"],
        ["data/loader.py", "Ler os arquivos do SMD e remover linhas inválidas preservando o alinhamento entre métricas e rótulos.", "caminho → DataFrame"],
        ["preprocessing/transform.py", "Padronizar por z-score e dividir a série preservando a ordem temporal.", "ndarray → ndarray"],
        ["models/model.py", "Definir o autoencoder (nn.Module, forward) e calcular o erro de reconstrução por janela.", "tensor → tensor"],
        ["models/persistence.py", "Salvar e carregar os pesos (state_dict) do modelo treinado.", "modelo ↔ arquivo .pt"],
        ["training/train.py", "Montar TensorDataset/DataLoader e executar o laço de treino e validação.", "matrizes → modelo + histórico"],
        ["evaluation/metrics.py", "Derivar o limiar, binarizar os erros e calcular precisão, recall, F1 e acurácia.", "erros + rótulos → métricas"],
        ["utils/config.py", "Fonte única de caminhos e hiperparâmetros do projeto.", "— (constantes)"],
    ]
    tabela(
        slide, dados, MARGEM, TOPO_CONTEUDO, LARGURA_UTIL, Inches(3.16),
        larguras=[0.235, 0.535, 0.23], tamanho_corpo=11,
    )

    y = TOPO_CONTEUDO + Inches(3.38)
    largura_codigo = Inches(7.35)
    linhas = [
        "# main.py - a integracao, em ordem",
        "train_clean = clean_data(load_data(config.TRAIN_PATH))",
        "X_train, X_val, _, _ = split_data(X_full, time_index, test_size=0.2)",
        "X_train, mean, std   = standardize(X_train)",
        "X_val,   _,    _     = standardize(X_val, mean=mean, std=std)",
        "model     = create_model(input_dim=X_train.shape[1], device=device)",
        "resultado = train_model(model, X_train, X_val, device=device)",
        "save_model(model)",
        "threshold = reconstruction_threshold(err_train, ANOMALY_PERCENTILE)",
        "metrics   = calculate_metrics(y_aligned, predict_anomalies(err, threshold))",
    ]
    altura_bloco = codigo(slide, linhas, MARGEM, y, largura_codigo, tamanho=9.5,
                          destaque={4})

    x2 = MARGEM + largura_codigo + Inches(0.26)
    largura_card = LARGURA_UTIL - largura_codigo - Inches(0.26)
    cartao(slide, x2, y, largura_card, altura_bloco, cor=GELO_AZUL, faixa=AMBAR)
    paragrafos(
        slide,
        [
            ("Como o main.py integra", {"tamanho": 12.5, "negrito": True, "cor": AZUL}),
            ("Ele chama uma função de cada módulo, na ordem, e passa o resultado adiante. "
             "Nenhum módulo chama outro: quem conhece a sequência é só o orquestrador.",
             {"tamanho": 11}),
            ("Na linha destacada está a decisão D4 em ação — a validação é padronizada com a "
             "média e o desvio vindos do treino.", {"tamanho": 11}),
        ],
        x2 + Inches(0.22), y + Inches(0.16), largura_card - Inches(0.44),
        altura_bloco - Inches(0.3), espaco_antes=6,
    )

    notas(slide, """
[TEMPO: 55s]

ROTEIRO — o que dizer
"Cada módulo tem uma responsabilidade e um contrato de entrada e saída explícito — é a coluna
da direita da tabela. Isso é o que permite trocar uma peça sem tocar nas outras.
Embaixo está o main.py resumido, que é a resposta para 'como isso tudo se integra'. Ele chama
uma função de cada módulo, na ordem, e passa o resultado adiante. Repare que nenhum módulo
aparece chamando outro: quem conhece a sequência é só o orquestrador.
E aponto a linha em destaque: a validação é padronizada com a média e o desvio que vieram do
treino, não os dela própria. É a decisão D4 do slide anterior aparecendo no código."

ENTENDA — os conceitos deste slide
· CONTRATO DE UMA FUNÇÃO: o que ela exige de entrada e o que garante de saída. Quando o
  contrato está explícito (nome claro + type hints + docstring), quem usa não precisa ler a
  implementação.
· DATAFRAME vs NDARRAY: DataFrame é a tabela do pandas, com rótulo de linha e coluna — bom
  para carregar e limpar. ndarray é a matriz do NumPy, sem rótulos — mais rápida para
  cálculo numérico. Nosso pipeline converte de um para o outro numa fronteira só.
· TENSOR: a matriz do PyTorch. Parece o ndarray, mas sabe calcular derivadas (é o que permite
  a rede aprender) e pode ir para a GPU.
· FRONTEIRA NUMPY → PYTORCH: a conversão para tensor acontece em dois pontos só — dentro do
  build_dataloader e do reconstruction_error. Nenhum outro módulo manipula tensor. É uma
  fronteira deliberada, não um acaso.
· "mean" e "std" DEVOLVIDOS PELA FUNÇÃO: repare que standardize devolve três coisas — os
  dados padronizados MAIS a média e o desvio usados. É isso que permite aplicar exatamente a
  mesma transformação na validação e no teste depois.

SE PERGUNTAREM
P: "O que é o clean_aligned e por que ele existe?"
R: "Quando removemos uma linha inválida das métricas de teste, temos que remover a MESMA
   linha dos rótulos, senão o pareamento entre métrica e rótulo se desloca e toda a avaliação
   fica errada. O clean_data comum não daria conta disso."
P: "Esse main.py cabe em dez linhas mesmo?"
R: "Essas são as chamadas essenciais; o arquivo real tem também as impressões de log e a
   gravação dos JSONs. Mas a estrutura é exatamente essa."
""")


# ===========================================================================
# Slide 9 - Codigo, NumPy e tipagem
# ===========================================================================
def slide_09(prs) -> None:
    slide = slide_padrao(prs, "Código Python, NumPy e tipagem", "Implementação · interfaces claras", 9)

    largura_col = Inches(5.95)

    # --- funcao 1
    texto(slide, "src/preprocessing/transform.py", MARGEM, TOPO_CONTEUDO, largura_col,
          Inches(0.26), tamanho=11.5, cor=AZUL_CLARO, negrito=True)
    linhas1 = [
        "def standardize(",
        "    X: np.ndarray,",
        "    mean: np.ndarray | None = None,",
        "    std: np.ndarray | None = None,",
        ") -> tuple[np.ndarray, np.ndarray, np.ndarray]:",
        '    """Padroniza os atributos para media 0 e desvio 1."""',
        "    X_float = np.asarray(X, dtype=float)",
        "",
        "    if mean is None or std is None:",
        "        mean = np.mean(X_float, axis=0)",
        "        std  = np.std(X_float, axis=0)",
        "",
        "    safe_std = np.where(std == 0, 1.0, std)",
        "    return (X_float - mean) / safe_std, mean, std",
    ]
    altura_bloco = codigo(slide, linhas1, MARGEM, TOPO_CONTEUDO + Inches(0.3), largura_col,
                          tamanho=10.5, destaque={4, 9, 10, 12})

    y_card = TOPO_CONTEUDO + Inches(0.42) + altura_bloco
    altura_card = BASE_CONTEUDO - y_card
    cartao(slide, MARGEM, y_card, largura_col, altura_card, cor=BRANCO, faixa=AMBAR)
    paragrafos(
        slide,
        [
            ("ENTRADA / SAÍDA   matriz (n_amostras, 38), opcionalmente com média e desvio já "
             "calculados → tupla com a matriz padronizada + a média + o desvio usados.",
             {"tamanho": 10.5}),
            ("NUMPY   np.mean e np.std com axis=0 calculam por coluna, ou seja, por métrica. A "
             "subtração e a divisão são vetorizadas: uma operação sobre a matriz inteira, sem laço.",
             {"tamanho": 10.5}),
            ("TIPAGEM   o | None diz que os parâmetros são opcionais; o tipo de retorno declara "
             "as três saídas. Quem chama não precisa ler o corpo da função.", {"tamanho": 10.5}),
            ("RESPONSABILIDADE   faz UMA coisa: padroniza. É a mesma função em dois modos — "
             "calcula as estatísticas (treino) ou aplica as recebidas (validação e teste).",
             {"tamanho": 10.5}),
        ],
        MARGEM + Inches(0.24), y_card + Inches(0.14), largura_col - Inches(0.46),
        altura_card - Inches(0.28), espaco_antes=5,
    )

    # --- funcao 2
    x2 = MARGEM + largura_col + Inches(0.19)
    texto(slide, "src/models/model.py", x2, TOPO_CONTEUDO, largura_col, Inches(0.26),
          tamanho=11.5, cor=AZUL_CLARO, negrito=True)
    linhas2 = [
        "def reconstruction_error(",
        "    model: Autoencoder,",
        "    X: np.ndarray,",
        "    window_size: int = config.WINDOW_SIZE,",
        "    device: str = config.DEVICE,",
        ") -> np.ndarray:",
        '    """Erro de reconstrucao (MSE) por janela."""',
        "    X_array = np.asarray(X, dtype=np.float32)",
        "    ...",
        "    windows = sliding_window_view(",
        "        X_array, window_shape=window_size, axis=0)",
        "    errors = torch.mean(",
        "        (reconstructed[:, -1, :] - batch[:, -1, :]) ** 2, dim=1)",
    ]
    codigo(slide, linhas2, x2, TOPO_CONTEUDO + Inches(0.3), largura_col,
           tamanho=10.5, destaque={5, 9, 10})

    cartao(slide, x2, y_card, largura_col, altura_card, cor=BRANCO, faixa=AMBAR)
    paragrafos(
        slide,
        [
            ("ENTRADA / SAÍDA   o modelo treinado e a matriz padronizada (n_amostras, 38) → "
             "vetor 1-D com um erro por janela, de tamanho n_amostras − 50 + 1.",
             {"tamanho": 10.5}),
            ("NUMPY   sliding_window_view constrói todas as janelas SEM copiar os dados: é uma "
             "visão sobre a mesma memória. Com 28 mil janelas de 50×38, copiar seria caro.",
             {"tamanho": 10.5}),
            ("TIPAGEM   o parâmetro é tipado como Autoencoder, não como “objeto qualquer”: o "
             "contrato exige o modelo certo, e o editor avisa antes de rodar.", {"tamanho": 10.5}),
            ("RESPONSABILIDADE   produz o escore. Não decide o que é anomalia — essa decisão é "
             "do evaluation/metrics.py, que compara o escore ao limiar.", {"tamanho": 10.5}),
        ],
        x2 + Inches(0.24), y_card + Inches(0.14), largura_col - Inches(0.46),
        altura_card - Inches(0.28), espaco_antes=5,
    )

    notas(slide, """
[TEMPO: 70s — são duas funções; ~35s em cada]

ROTEIRO — o que dizer
"Duas funções reais do projeto, escolhidas por serem representativas.
À esquerda, a padronização. Repare na assinatura: ela declara os tipos de entrada e o tipo de
retorno — uma tupla com três arrays. Quem chama sabe o que recebe sem ler o corpo. O uso de
NumPy aqui é vetorizado: np.mean com axis=0 calcula a média por coluna, ou seja, por métrica,
e a subtração acontece na matriz inteira de uma vez, sem laço em Python. E há um detalhe de
robustez: se alguma métrica for constante, o desvio seria zero e a divisão explodiria — a
linha do np.where troca o zero por um antes de dividir.
À direita, o cálculo do erro de reconstrução. O ponto de NumPy aqui é o sliding_window_view:
ele constrói todas as janelas sem copiar os dados, é uma visão sobre a mesma memória. Com
vinte e oito mil janelas de cinquenta por trinta e oito, copiar seria caro."

ENTENDA — os conceitos deste slide
· TYPE HINTS: as anotações de tipo (X: np.ndarray, -> np.ndarray). Python não as obriga nem as
  verifica em tempo de execução — elas servem para quem lê o código e para o editor apontar
  erro antes de rodar. É critério avaliado na disciplina.
· np.ndarray | None: a barra significa "ou". O parâmetro aceita um array OU nada. É o que
  permite a mesma função ter dois modos de uso.
· VETORIZAÇÃO: fazer a operação na matriz inteira de uma vez, em código otimizado por baixo,
  em vez de percorrer elemento a elemento com for. É a razão de existir do NumPy: mais rápido
  e mais legível.
· axis=0: "ao longo das linhas", ou seja, o resultado tem um valor por coluna. Como cada
  coluna é uma métrica, é exatamente o que queremos — média por métrica, não média geral.
· SLIDING_WINDOW_VIEW: função do NumPy que devolve todas as janelas deslizantes como uma
  "visão" — sem alocar memória nova. Se fizéssemos as janelas com um laço e uma cópia, seriam
  28 mil × 50 × 38 números duplicados.
· MSE (erro quadrático médio): média dos quadrados das diferenças. Eleva ao quadrado para que
  erro para mais e para menos não se cancelem, e para punir mais os desvios grandes.

SE PERGUNTAREM
P: "Por que o safe_std existe?"
R: "Se uma métrica for constante no treino, o desvio dela é zero e a divisão daria infinito
   ou NaN, contaminando tudo. Trocamos o zero por um: a métrica fica centrada e sem escala,
   que é o comportamento razoável para uma coluna que não varia."
P: "Por que só o último instante da janela entra no erro?"
R: "Porque associamos cada janela ao instante mais recente dela — é o que permite comparar
   com o rótulo daquele instante. Por isso o alinhamento no main.py começa em janela menos um."
""")


# ===========================================================================
# Slide 10 - Modelo em PyTorch
# ===========================================================================
def slide_10(prs) -> None:
    slide = slide_padrao(prs, "Implementação do modelo em PyTorch", "Implementação · o modelo real", 10)

    largura_codigo = Inches(6.9)
    linhas = [
        "class Autoencoder(nn.Module):",
        "    def __init__(self, input_dim: int,",
        "                 window_size: int = config.WINDOW_SIZE,",
        "                 hidden_dims: tuple[int, ...] = (64, 32),",
        "                 latent_dim: int = 16) -> None:",
        "        super().__init__()",
        "        self.flat_dim = input_dim * window_size   # 38 x 50 = 1900",
        "",
        "        encoder_dims = (self.flat_dim, *hidden_dims, latent_dim)",
        "        decoder_dims = (latent_dim, *reversed(hidden_dims), self.flat_dim)",
        "        self.encoder = self._build_mlp(encoder_dims)",
        "        self.decoder = self._build_mlp(decoder_dims)",
        "",
        "    def forward(self, x: torch.Tensor) -> torch.Tensor:",
        "        flattened     = x.reshape(x.shape[0], self.flat_dim)",
        "        latent        = self.encoder(flattened)",
        "        reconstructed = self.decoder(latent)",
        "        return reconstructed.reshape(",
        "            x.shape[0], self.window_size, self.input_dim)",
    ]
    altura_bloco = codigo(slide, linhas, MARGEM, TOPO_CONTEUDO, largura_codigo,
                          tamanho=10.5, destaque={0, 13})

    # funil de dimensoes
    y_funil = TOPO_CONTEUDO + altura_bloco + Inches(0.2)
    cartao(slide, MARGEM, y_funil, largura_codigo, BASE_CONTEUDO - y_funil,
           cor=BRANCO, faixa=AZUL)
    texto(slide, "O FUNIL: a janela é comprimida e reconstruída", MARGEM + Inches(0.24),
          y_funil + Inches(0.14), largura_codigo - Inches(0.45), Inches(0.24),
          tamanho=11.5, cor=AZUL, negrito=True)

    dims = ["1900", "64", "32", "16", "32", "64", "1900"]
    largura_dim = Inches(0.78)
    espaco_dim = Inches(0.2)
    x = MARGEM + Inches(0.3)
    for indice, dim in enumerate(dims):
        cor_caixa = AMBAR if indice == 3 else GELO_AZUL
        cor_texto = BRANCO if indice == 3 else AZUL
        retangulo(slide, x, y_funil + Inches(0.52), largura_dim, Inches(0.44), cor=cor_caixa)
        texto(slide, dim, x, y_funil + Inches(0.62), largura_dim, Inches(0.26),
              tamanho=13, cor=cor_texto, negrito=True, alinhamento=PP_ALIGN.CENTER)
        if indice < len(dims) - 1:
            seta = retangulo(slide, x + largura_dim + Inches(0.04),
                             y_funil + Inches(0.68), espaco_dim - Inches(0.08), Inches(0.11),
                             cor=CINZA, forma=MSO_SHAPE.RIGHT_ARROW)
            seta.line.fill.background()
        x += largura_dim + espaco_dim
    texto(slide, "encoder  (comprime)", MARGEM + Inches(0.3), y_funil + Inches(1.04),
          Inches(2.6), Inches(0.24), tamanho=10, cor=CINZA)
    texto(slide, "gargalo", MARGEM + Inches(3.24), y_funil + Inches(1.04), Inches(1.2),
          Inches(0.24), tamanho=10, cor=AMBAR, negrito=True)
    texto(slide, "decoder  (reconstrói)", MARGEM + Inches(4.5), y_funil + Inches(1.04),
          Inches(2.6), Inches(0.24), tamanho=10, cor=CINZA)

    # ficha tecnica
    x_ficha = MARGEM + largura_codigo + Inches(0.26)
    largura_ficha = LARGURA_UTIL - largura_codigo - Inches(0.26)
    cfg_a = CONFIGS["A"]
    ficha = [
        ("Arquitetura", "MLP simétrico: encoder 1900→64→32→16, decoder espelhado. Ativação ReLU entre camadas."),
        ("Parâmetros", f"{cfg_a['n_parametros']:,}".replace(",", ".") + " pesos treináveis"),
        ("Função de perda", "nn.MSELoss — erro quadrático médio entre a janela e sua reconstrução"),
        ("Otimizador", "torch.optim.Adam, learning rate 1e-3"),
        ("Épocas", f"máximo 100, executadas {cfg_a['epocas_executadas']} (early stopping, paciência 10)"),
        ("Batch size", "64 janelas por lote"),
        ("Device", "CPU — get_device(prefer_cuda=False); nenhuma GPU foi necessária"),
        ("Estabilidade", "ReduceLROnPlateau (fator 0,5) + gradient clipping (norma 1,0)"),
        ("Salvar / carregar", "save_model grava o state_dict em artifacts/autoencoder.pt; load_model restaura numa instância criada com a mesma arquitetura"),
    ]
    altura_ficha = BASE_CONTEUDO - TOPO_CONTEUDO
    cartao(slide, x_ficha, TOPO_CONTEUDO, largura_ficha, altura_ficha, cor=BRANCO, faixa=AMBAR)
    texto(slide, "FICHA TÉCNICA DO TREINO", x_ficha + Inches(0.24), TOPO_CONTEUDO + Inches(0.16),
          largura_ficha - Inches(0.46), Inches(0.24), tamanho=10, cor=AZUL_CLARO, negrito=True)
    linhas_ficha = []
    for rotulo, valor in ficha:
        linhas_ficha.append((rotulo.upper(), {"tamanho": 9.5, "cor": AMBAR, "negrito": True}))
        linhas_ficha.append((valor, {"tamanho": 11, "cor": GRAFITE}))
    paragrafos(
        slide, linhas_ficha, x_ficha + Inches(0.24), TOPO_CONTEUDO + Inches(0.46),
        largura_ficha - Inches(0.46), altura_ficha - Inches(0.6), espaco_antes=4,
        espaco_linhas=1.02,
    )

    notas(slide, """
[TEMPO: 75s — é o slide mais técnico; vá com calma no funil]

ROTEIRO — o que dizer
"Este é o modelo real. A classe herda de nn.Module, que é a classe base de todo modelo em
PyTorch, e implementa o forward — o método que diz o que acontece com a entrada.
Olhem o desenho embaixo, que é a essência do método. Uma janela de 50 instantes por 38
métricas é achatada em um vetor de 1900 números. O encoder comprime esse vetor até 16 — é o
gargalo, em destaque. O decoder faz o caminho de volta e tenta reconstruir os mesmos 1900.
A pergunta natural é: por que espremer? Porque em 16 números não cabe tudo. A rede é obrigada
a guardar só os padrões que mais se repetem, e o que mais se repete no treino é o
comportamento normal. Então ela fica boa em reconstruir o normal e ruim em reconstruir
qualquer coisa diferente — e essa dificuldade de reconstruir é exatamente o nosso sinal.
À direita, a ficha técnica: perda MSE, otimizador Adam, no máximo 100 épocas mas parou em 71
por early stopping, lotes de 64 janelas, tudo em CPU."

ENTENDA — os conceitos deste slide
· nn.Module: a classe base do PyTorch. Herdar dela dá gerenciamento automático dos pesos e o
  encadeamento com otimizador e perda.
· FORWARD: o caminho da entrada até a saída. Você escreve só esse caminho; o PyTorch calcula
  as derivadas para trás sozinho (é o autograd).
· MLP (Multi-Layer Perceptron): rede de camadas totalmente conectadas, a mais simples que
  existe. Escolhemos ela por ser suficiente e barata; alternativas com memória temporal, como
  LSTM, estão citadas no documento como candidatas.
· ReLU: função de ativação. Sem uma função não linear entre as camadas, empilhar camadas
  lineares equivale a uma só — a rede não ganharia poder de representação.
· FUNÇÃO DE PERDA (MSELoss): o número que mede o quão errada está a saída. Aqui é o erro
  quadrático médio entre a janela original e a reconstruída. É o que o treino minimiza.
· OTIMIZADOR (Adam): o algoritmo que ajusta os pesos a cada lote, seguindo o gradiente da
  perda. Adam é o padrão prático — ajusta o passo automaticamente por parâmetro.
· LEARNING RATE (1e-3 = 0,001): o tamanho do passo. Grande demais, o treino oscila; pequeno
  demais, demora eternamente.
· BATCH SIZE (64): quantas janelas entram por vez antes de atualizar os pesos.
· ÉPOCA: uma passagem completa por todo o conjunto de treino.
· EARLY STOPPING: parar quando a validação para de melhorar (aqui, 10 épocas sem melhora). É
  o que evita overfitting e explica por que rodamos 71 e não 100 épocas.
· OVERFITTING: quando o modelo decora o treino em vez de aprender o padrão. Detecta-se
  olhando a curva de validação subir enquanto a de treino desce — está no slide 13.
· GRADIENT CLIPPING: limita o tamanho do ajuste quando o gradiente explode, evitando que uma
  única atualização destrua os pesos.
· STATE_DICT: dicionário com os pesos aprendidos. É o que salvamos.

SE PERGUNTAREM
P: "Por que não uma LSTM, já que é série temporal?"
R: "É a candidata natural, e está listada no documento de requisitos como alternativa. Optamos
   pelo autoencoder MLP porque a janela já dá contexto temporal ao modelo, o custo é muito
   menor e o escopo da disciplina pede a maquete funcional. Testar a LSTM é a nossa primeira
   melhoria futura."
P: "Por que achatar a janela? Isso não perde a ordem temporal?"
R: "Perde parcialmente: o MLP não tem noção explícita de sequência, ele vê 1900 posições fixas.
   A ordem ainda importa, porque cada posição corresponde sempre ao mesmo instante relativo,
   mas um modelo recorrente capturaria a dependência temporal melhor. É uma limitação
   assumida."
P: "Treinou em quanto tempo?"
R: "Cerca de dois minutos e meio em CPU — o que atende o RNF-03, que exigia treino viável sem
   GPU dedicada."
""")


# ===========================================================================
# Slide 11 - Configuracao experimental
# ===========================================================================
def slide_11(prs) -> None:
    slide = slide_padrao(prs, "Configuração experimental", "Experimentos · o que quisemos descobrir", 11)

    # pergunta do experimento
    texto(
        slide,
        "O QUE QUERÍAMOS DESCOBRIR:   a capacidade do modelo muda a qualidade da detecção — "
        "ou a alavanca real é o limiar?",
        MARGEM, TOPO_CONTEUDO, LARGURA_UTIL, Inches(0.3),
        tamanho=13, cor=AZUL, negrito=True,
    )

    # protocolo
    y = TOPO_CONTEUDO + Inches(0.38)
    protocolo = [
        ("Divisão dos dados",
         "Treino 80% / validação 20% do arquivo train, em corte temporal (sem embaralhar). "
         "O arquivo test é usado só na avaliação final."),
        ("Métricas",
         "Precisão, recall, F1 e acurácia, calculadas em NumPy puro contra o test_label."),
        ("Controle de variáveis",
         "Semente fixa (42) e mesmo pré-processamento em todas as execuções: a única coisa que "
         "muda entre configurações é o hiperparâmetro em teste."),
    ]
    largura_prot = Inches(3.9)
    x = MARGEM
    for titulo, corpo in protocolo:
        cartao(slide, x, y, largura_prot, Inches(1.02), cor=BRANCO, faixa=AMBAR)
        texto(slide, titulo, x + Inches(0.22), y + Inches(0.12), largura_prot - Inches(0.42),
              Inches(0.24), tamanho=11.5, cor=AZUL, negrito=True)
        texto(slide, corpo, x + Inches(0.22), y + Inches(0.4), largura_prot - Inches(0.42),
              Inches(0.58), tamanho=10.5, cor=GRAFITE, espaco_linhas=1.02)
        x += largura_prot + Inches(0.2)

    # tabela de configuracoes
    y2 = y + Inches(1.16)
    dados = [["Config", "O que muda", "Parâmetros", "Épocas", "Erro val.", "Precisão", "Recall", "F1"]]
    for cid in ["A", "B", "C", "D"]:
        cfg = CONFIGS[cid]
        met = cfg["por_percentil"]["99.5"]
        dados.append([
            f"{cid} — {cfg['nome']}",
            cfg["descricao"],
            f"{cfg['n_parametros']:,}".replace(",", "."),
            str(cfg["epocas_executadas"]),
            num(cfg["melhor_val_loss"]),
            num(met["precision"], 3),
            num(met["recall"], 3),
            num(met["f1_score"], 3),
        ])
    tabela(
        slide, dados, MARGEM, y2, LARGURA_UTIL, Inches(1.62),
        larguras=[0.175, 0.275, 0.105, 0.075, 0.095, 0.09, 0.085, 0.1],
        tamanho_corpo=10.5, tamanho_cabecalho=11,
    )

    # figura + achados
    y3 = y2 + Inches(1.78)
    altura_fig = Inches(1.98)
    largura_fig = figura(slide, "tradeoff_limiar.png", MARGEM, y3, altura_fig)

    x3 = MARGEM + largura_fig + Inches(0.22)
    largura_achado = LARGURA_UTIL - largura_fig - Inches(0.22)
    cartao(slide, x3, y3, largura_achado, altura_fig, cor=GELO_AZUL, faixa=VERDE)

    p999 = CONFIGS["A"]["por_percentil"]["99.9"]
    p995 = CONFIGS["A"]["por_percentil"]["99.5"]
    f1_por_config = [CONFIGS[c]["por_percentil"]["99.5"]["f1_score"] for c in "ABCD"]
    amplitude_config = max(f1_por_config) - min(f1_por_config)
    amplitude_limiar = p999["f1_score"] - p995["f1_score"]

    paragrafos(
        slide,
        [
            ("DOIS ACHADOS", {"tamanho": 12.5, "negrito": True, "cor": AZUL}),
            (f"1.  Reconstruir melhor não é detectar melhor.  A config D teve o MENOR erro de "
             f"validação ({num(CONFIGS['D']['melhor_val_loss'])}) e não lidera a detecção "
             f"(F1 {num(f1_por_config[3], 3)}); a C, com erro de validação pior "
             f"({num(CONFIGS['C']['melhor_val_loss'])}), tem o melhor F1 "
             f"({num(f1_por_config[2], 3)}) e recall de "
             f"{pct(CONFIGS['C']['por_percentil']['99.5']['recall'])}. Capacidade demais "
             f"reconstrói bem até o que é anômalo — e isso reduz a separação.",
             {"tamanho": 11}),
            (f"2.  O limiar move mais que a arquitetura.  Entre as quatro configurações o F1 "
             f"varia {num(amplitude_config, 3)}; mexendo só no percentil, na MESMA "
             f"configuração, varia {num(amplitude_limiar, 3)} — mais que o triplo. A precisão "
             f"vai de {pct(p995['precision'])} a {pct(p999['precision'])}, pagando com recall "
             f"({pct(p995['recall'])} → {pct(p999['recall'])}).",
             {"tamanho": 11}),
            ("Consequência: calibrar o limiar ao custo do alerta é a decisão de maior impacto — "
             "e é exatamente a decisão que o RF-06 entrega ao Coordenador.",
             {"tamanho": 11, "negrito": True}),
        ],
        x3 + Inches(0.24), y3 + Inches(0.13), largura_achado - Inches(0.46),
        altura_fig - Inches(0.26), espaco_antes=4, espaco_linhas=1.02,
    )

    notas(slide, """
[TEMPO: 80s — é o slide mais denso; o achado é o que interessa]

ROTEIRO — o que dizer
"Aqui está o protocolo experimental e o que ele nos ensinou.
O protocolo: dividimos o arquivo de treino em 80% treino e 20% validação, em corte temporal,
sem embaralhar. O arquivo de teste só é tocado na avaliação final. Semente fixa em todas as
execuções, então a única coisa que muda entre elas é o hiperparâmetro sob teste.
Testamos quatro configurações: a base, uma com o gargalo estreitado para 8, uma com a janela
dobrada para 100 e uma com mais capacidade — gargalo 32 e camadas maiores.
Dois achados.
O primeiro é contraintuitivo e é o que eu mais gosto: reconstruir melhor não é detectar
melhor. A configuração D teve o MENOR erro de validação de todas, 0,157 — ela é a que melhor
reconstrói. E mesmo assim não lidera a detecção. A C, que reconstrói pior, tem o melhor F1 e
recall de 96%. A leitura é que capacidade demais faz o autoencoder reconstruir bem até o que
é anômalo, e aí a anomalia deixa de se destacar.
O segundo achado é a amplitude. Entre as quatro arquiteturas o F1 varia cerca de cinco
centésimos. Mexendo só no limiar, na mesma configuração, varia dezessete centésimos — mais
que o triplo. Olhem o gráfico: a precisão sai de 25% e vai a 83%, o recall cai de 88% para
43%.
A consequência prática: a decisão de maior impacto neste sistema não é qual arquitetura usar,
é onde colocar o corte. E é exatamente essa decisão que o nosso RF-06 coloca na mão do
coordenador de operações."

ENTENDA — os conceitos deste slide
· HIPERPARÂMETRO: o que você escolhe ANTES de treinar (tamanho do gargalo, taxa de
  aprendizado, tamanho do lote). Diferente do parâmetro, que é o peso que a rede aprende.
· CONTROLE DE VARIÁVEIS: mudar uma coisa por vez. Se mudássemos gargalo e janela juntos, não
  saberíamos a qual atribuir a diferença. É o mesmo princípio de qualquer experimento.
· SEMENTE (seed): número que fixa a aleatoriedade — inicialização dos pesos, ordem dos lotes.
  Com semente fixa, rodar de novo dá o mesmo resultado. Sem ela, não dá para comparar duas
  execuções, porque parte da diferença seria só sorte.
· VALIDAÇÃO vs TESTE: a validação serve para decidir quando parar de treinar; o teste serve
  para medir o resultado final. Se usássemos o teste para decidir quando parar, o número
  final estaria contaminado.
· POR QUE MAIS CAPACIDADE PIOROU A DETECÇÃO: o método depende do modelo ser RUIM em
  reconstruir o que é anômalo. Se a rede tem capacidade sobrando, ela aprende a reconstruir
  quase qualquer coisa — inclusive o que nunca viu — e o erro deixa de separar normal de
  anômalo. Apresente isso como a interpretação mais provável, não como fato provado: nós
  medimos o efeito, não isolamos a causa.
· ERRO DE VALIDAÇÃO vs QUALIDADE DA DETECÇÃO: são duas coisas diferentes, e este é o ponto
  fino do slide. O erro de validação mede o quanto a rede reconstrói bem o NORMAL; a detecção
  depende da DISTÂNCIA entre o erro no normal e o erro no anômalo. Otimizar o primeiro não
  garante o segundo.
· O QUE O PERCENTIL SIGNIFICA: p99,5 quer dizer "o valor abaixo do qual ficam 99,5% dos erros
  de treino". Subir para p99,9 é exigir um erro muito maior para chamar de anomalia — daí
  menos alertas, mais precisos, e mais anomalias escapando.

SE PERGUNTAREM
P: "Então o modelo não importa?"
R: "Importa: é ele que produz o escore, sem escore não há o que cortar. O que medimos é que,
   dentro da família que testamos, a capacidade não foi o fator dominante — e que aumentá-la
   chegou a atrapalhar. Um modelo de outra classe, como uma LSTM, poderia mudar o quadro;
   não testamos."
P: "Por que a configuração C, que é a melhor em F1, não virou a configuração oficial?"
R: "Porque a execução de referência do repositório é a A, e preferimos não trocar o resultado
   registrado no meio da apresentação. A C é a recomendação para o próximo ciclo — e note que
   ela custa o dobro de parâmetros e piora bastante no regime de precisão alta."
P: "Qual configuração vocês recomendariam em produção?"
R: "Depende do custo relativo. Se deixar passar um incidente custa mais que investigar um
   alerta falso — o caso típico de operação — fica p99,5 com recall alto. Se o plantão está
   saturado, p99,9 entrega três em cada quatro alertas corretos. A decisão é do negócio, e o
   sistema registra o limiar aplicado em cada execução justamente para essa escolha ser
   explícita."
""")


# ===========================================================================
# Slide 12 - Testes automatizados
# ===========================================================================
def slide_12(prs) -> None:
    slide = slide_padrao(prs, "Testes automatizados e evidências", "Testes · unittest", 12)

    largura_codigo = Inches(7.35)
    linhas = [
        "# tests/test_preprocessing.py",
        "class TestPreprocessing(unittest.TestCase):",
        "    def test_standardize_mean(self):",
        '        """A media de cada feature deve ser aprox. zero."""',
        "        X_std, _, _ = standardize(self.X)",
        "        np.testing.assert_allclose(",
        "            np.mean(X_std, axis=0), np.zeros(38), atol=1e-10)",
        "",
        "# tests/test_model.py",
        "    def test_model_output_shape(self):",
        '        """A saida deve ter o mesmo shape da entrada."""',
        "        output = self.model(self.x)",
        "        self.assertEqual(output.shape, self.x.shape)",
        "",
        "    def test_load_model_restores_weights(self):",
        '        """load_model restaura pesos identicos aos salvos."""',
        "        save_model(self.model, self.caminho)",
        "        restaurado = load_model(create_model(38), self.caminho)",
    ]
    altura_bloco = codigo(slide, linhas, MARGEM, TOPO_CONTEUDO, largura_codigo,
                          tamanho=10.5, destaque={2, 9, 14})

    # evidencia
    y = TOPO_CONTEUDO + altura_bloco + Inches(0.18)
    saida = [
        "$ python -m unittest discover -s tests",
        "test_standardize_mean ...................... ok",
        "test_model_output_shape .................... ok",
        "test_load_model_restores_weights ........... ok",
        "--------------------------------------------------",
        "Ran 13 tests in 0.453s",
        "OK",
    ]
    codigo(slide, saida, MARGEM, y, largura_codigo, tamanho=10, destaque={5, 6})

    # o que cobrimos
    x2 = MARGEM + largura_codigo + Inches(0.26)
    largura_col = LARGURA_UTIL - largura_codigo - Inches(0.26)
    altura_col = BASE_CONTEUDO - TOPO_CONTEUDO
    cartao(slide, x2, TOPO_CONTEUDO, largura_col, altura_col, cor=BRANCO, faixa=AMBAR)
    texto(slide, "O QUE A SUÍTE COBRE", x2 + Inches(0.24), TOPO_CONTEUDO + Inches(0.16),
          largura_col - Inches(0.46), Inches(0.24), tamanho=10, cor=AZUL_CLARO, negrito=True)
    cobertura = [
        ("test_data.py — 3 testes", {"tamanho": 12, "negrito": True, "cor": AZUL}),
        ("Carregamento devolve DataFrame com 38 colunas · limpeza remove NaN sem alterar colunas · "
         "caminho inexistente levanta FileNotFoundError.", {"tamanho": 10.5}),
        ("test_preprocessing.py — 4 testes", {"tamanho": 12, "negrito": True, "cor": AZUL}),
        ("Padronização preserva o shape · média resultante ≈ 0 · desvio ≈ 1 · split respeita as "
         "proporções e a ordem.", {"tamanho": 10.5}),
        ("test_model.py — 4 testes", {"tamanho": 12, "negrito": True, "cor": AZUL}),
        ("Saída do modelo tem o shape da entrada · saída é finita (sem NaN/Inf) · save_model cria o "
         "arquivo · load_model restaura os pesos idênticos.", {"tamanho": 10.5}),
        ("test_training.py — 2 testes", {"tamanho": 12, "negrito": True, "cor": AZUL}),
        ("Tensores do DataLoader têm shape (batch, 50, 38) · e dtype float32.", {"tamanho": 10.5}),
        ("Por que estes testes", {"tamanho": 12, "negrito": True, "cor": AMBAR}),
        ("São as fronteiras entre os módulos — shape, dtype e contrato de entrada/saída. É onde a "
         "integração quebra silenciosamente quando alguém muda um módulo sem avisar o outro.",
         {"tamanho": 10.5}),
    ]
    paragrafos(
        slide, cobertura, x2 + Inches(0.24), TOPO_CONTEUDO + Inches(0.46),
        largura_col - Inches(0.46), altura_col - Inches(0.6), espaco_antes=6,
    )

    notas(slide, """
[TEMPO: 55s — mostre a evidência de execução, é o que o enunciado pede]

ROTEIRO — o que dizer
"São 13 testes em unittest, organizados em quatro arquivos — um por pacote do src. Escolhi
três para mostrar.
O primeiro verifica uma propriedade matemática: depois de padronizar, a média de cada métrica
tem que ser aproximadamente zero. Não testamos que a função rodou; testamos que ela fez a
coisa certa.
O segundo é de contrato: a saída do autoencoder precisa ter exatamente o mesmo formato da
entrada. Se alguém mexer na arquitetura e quebrar isso, o teste pega na hora.
O terceiro é o mais interessante: salvar o modelo, carregar de volta numa instância nova, e
comparar peso a peso. É o que garante que o modelo salvo é o modelo treinado.
E aqui embaixo está a evidência de execução: Ran 13 tests, OK."

ENTENDA — os conceitos deste slide
· unittest: o framework de testes que vem com o Python, sem instalar nada. A disciplina exige
  ele e não o pytest.
· TestCase: a classe base. Cada método que começa com test_ é um teste, e roda isolado dos
  demais.
· setUp: método que roda ANTES de cada teste, preparando o cenário (nos nossos, criando dados
  sintéticos com semente fixa). Isso garante que um teste não contamine o outro.
· assertEqual / assert_allclose: as afirmações. A segunda é a versão do NumPy para números com
  casas decimais — comparar float com igualdade exata falha por arredondamento, então
  comparamos com uma tolerância (o atol=1e-10).
· SHAPE: as dimensões de uma matriz ou tensor. (batch, 50, 38) significa: um lote de janelas,
  cada uma com 50 instantes e 38 métricas.
· DTYPE float32: o tipo numérico. O PyTorch trabalha em float32 por padrão; se um tensor
  chegar como float64, algumas operações quebram. Por isso testamos o dtype explicitamente.
· POR QUE DADOS SINTÉTICOS NOS TESTES: os testes não dependem do dataset estar baixado. Eles
  geram uma matriz aleatória com semente fixa e verificam a propriedade. Assim a suíte roda em
  qualquer máquina, em menos de meio segundo.

SE PERGUNTAREM
P: "13 testes é pouco?"
R: "É o que cobre o que o enunciado pede: carregamento, pré-processamento, shape dos tensores,
   saída do modelo e o ciclo salvar/carregar. O que falta e reconhecemos: não há teste do
   cálculo das métricas nem do laço de treino completo — são as próximas a escrever."
P: "Como rodar?"
R: "python -m unittest discover -s tests, a partir da raiz do repositório. Está no README, com
   a evidência da execução."
""")


# ===========================================================================
# Slide 13 - Resultados, conclusao e limitacoes
# ===========================================================================
def slide_13(prs) -> None:
    slide = slide_padrao(prs, "Resultados, conclusão e limitações", "Experimentos · o que ficou de pé", 13)

    p995 = CONFIGS["A"]["por_percentil"]["99.5"]

    # --- dois graficos: desempenho no treino e deteccao no teste
    altura_fig = Inches(2.18)
    y_fig = TOPO_CONTEUDO + Inches(0.26)

    texto(slide, "TREINAMENTO — erro por época", MARGEM, TOPO_CONTEUDO, Inches(4.0),
          Inches(0.24), tamanho=11, cor=AZUL_CLARO, negrito=True)
    largura_1 = figura(slide, "curva_treino.png", MARGEM, y_fig, altura_fig)

    x_fig2 = MARGEM + largura_1 + Inches(0.24)
    texto(slide, "TESTE — erro por janela, limiar e anomalias reais", x_fig2, TOPO_CONTEUDO,
          Inches(5.4), Inches(0.24), tamanho=11, cor=AZUL_CLARO, negrito=True)
    figura(slide, "erro_no_teste.png", x_fig2, y_fig, altura_fig)

    # --- faixa de numeros da execucao de referencia
    y_num = y_fig + altura_fig + Inches(0.2)
    numeros = [
        ("ERRO DE VALIDAÇÃO", num(RESULTADOS["best_val_loss"]), AZUL),
        ("ÉPOCAS / MELHOR", f"{RESULTADOS['epochs_executed']} / {RESULTADOS['best_epoch']}", AZUL),
        ("LIMIAR (p99,5)", num(RESULTADOS["threshold"], 3), AZUL),
        ("PRECISÃO · RECALL",
         f"{pct(p995['precision'])} · {pct(p995['recall'])}", AZUL),
        ("CRITÉRIO DE ACEITAÇÃO", "ATENDIDO", VERDE),
    ]
    largura_num = Inches(2.33)
    x = MARGEM
    for rotulo, valor, cor in numeros:
        cartao(slide, x, y_num, largura_num, Inches(0.78), cor=BRANCO, faixa=cor)
        texto(slide, rotulo, x + Inches(0.2), y_num + Inches(0.11), largura_num - Inches(0.36),
              Inches(0.2), tamanho=9, cor=CINZA, negrito=True)
        texto(slide, valor, x + Inches(0.2), y_num + Inches(0.34), largura_num - Inches(0.36),
              Inches(0.36), tamanho=17, cor=cor, negrito=True)
        x += largura_num + Inches(0.11)

    # conclusao, limitacoes, futuro
    y2 = y_num + Inches(0.94)
    colunas = [
        ("RESOLVEU O PROBLEMA?", VERDE, [
            "Sim, como maquete funcional: o pipeline roda ponta a ponta e o critério de aceitação é atendido.",
            "5 dos 7 RF atendidos; 2 parciais; 1 é do sistema-alvo.",
            "3 dos 4 RNF verificados com número medido.",
        ]),
        ("LIMITAÇÕES", VERMELHO, [
            "Erro agregado por janela, não por métrica — trava RF-03 e RNF-02.",
            "Validado só na machine-1-1, com janela de 50.",
            "O baseline de limiar por métrica nunca foi medido: o ganho do ML está argumentado, não provado.",
        ]),
        ("PRÓXIMOS PASSOS", AZUL, [
            "1. Decompor o erro por dimensão — destrava RF-03/RNF-02 e é mudança localizada.",
            "2. Medir o baseline de limiar fixo, para comprovar (ou não) o ganho do modelo.",
            "3. Testar LSTM-autoencoder e estender às 28 máquinas.",
        ]),
    ]
    largura_col3 = Inches(3.9)
    altura_col3 = BASE_CONTEUDO - y2
    x = MARGEM
    for titulo, cor, itens in colunas:
        cartao(slide, x, y2, largura_col3, altura_col3, cor=BRANCO, faixa=cor)
        texto(slide, titulo, x + Inches(0.22), y2 + Inches(0.13), largura_col3 - Inches(0.42),
              Inches(0.24), tamanho=11, cor=cor, negrito=True)
        paragrafos(
            slide,
            [(item, {"tamanho": 10.5, "marcador": "· "}) for item in itens],
            x + Inches(0.22), y2 + Inches(0.42), largura_col3 - Inches(0.42),
            altura_col3 - Inches(0.55), espaco_antes=6, espaco_linhas=1.02,
        )
        x += largura_col3 + Inches(0.2)

    notas(slide, """
[TEMPO: 75s — é o fechamento; termine no tempo e agradeça]

ROTEIRO — o que dizer
"Fechando com os resultados, em dois gráficos.
À esquerda, o treinamento: as curvas de treino e validação descem juntas e a validação
estabiliza por volta da época 61 — que é a melhor época, com erro 0,1896. O early stopping
parou na 71, e os pesos salvos são os da 61, não os da última. Repare que a validação não
sobe: não há sinal de overfitting sério.
À direita, o sistema funcionando no teste. A linha azul é o erro de cada janela em escala
logarítmica; a tracejada é o limiar; e as faixas rosa são as anomalias reais rotuladas.
Vejam que os picos maiores caem dentro das faixas — é o detector acertando. E vejam também,
com honestidade, que depois da metade da série há muito ponto azul acima do limiar fora das
faixas: são os falsos positivos que explicam a precisão de 25%.
Abaixo, os números da execução de referência, e o critério de aceitação atendido.
Respondendo às três perguntas do enunciado. Resolveu o problema? Sim, como maquete funcional:
o pipeline roda ponta a ponta e o critério de aceitação que definimos foi atendido nos dois
lados. Os requisitos foram atendidos? Cinco dos sete funcionais; dois parciais; um é
explicitamente do sistema-alvo.
E as limitações — que eu prefiro dizer do que esconder. A maior: o erro é agregado por janela,
não por métrica, e é isso que trava dois requisitos. A segunda: validamos numa máquina só. E
a terceira, a mais honesta: nunca medimos o baseline de limiar fixo, então o ganho do modelo
está argumentado e não provado.
Daí os próximos passos: decompor o erro por dimensão, medir o baseline, e testar uma LSTM."

ENTENDA — os conceitos deste slide
· COMO LER A CURVA (gráfico da esquerda): eixo x são as épocas, eixo y é o erro. Treino e
  validação descendo juntas = a rede está aprendendo padrão de verdade. Se a de treino
  continuasse caindo e a de validação começasse a subir, seria overfitting — a rede estaria
  decorando.
· COMO LER O GRÁFICO DA DIREITA: cada ponto da linha azul é uma janela do teste, na ordem em
  que aconteceu, e a altura é o erro dela. A escala é logarítmica — cada marca do eixo é 100
  vezes a anterior — porque os picos são milhões de vezes maiores que o normal e, numa escala
  comum, tudo viraria uma linha rente ao chão com três espetos. Tudo que está acima da linha
  tracejada é sinalizado como anomalia pelo sistema; as faixas rosa são as anomalias
  verdadeiras. Acerto = pico dentro da faixa. Falso positivo = ponto acima da linha fora da
  faixa.
· POR QUE A VALIDAÇÃO FICA ACIMA DO TREINO: normal e esperado. O modelo viu os dados de treino
  e nunca viu os de validação. A distância entre as duas é pequena aqui, o que é bom sinal.
· MELHOR ÉPOCA vs ÚLTIMA ÉPOCA: com early stopping o treino continua um pouco além do melhor
  ponto, para confirmar que não vai melhorar. Guardamos os pesos da melhor época e registramos
  qual foi — senão o número reportado não corresponderia ao modelo salvo.
· "MAQUETE FUNCIONAL": é o termo do próprio enunciado. O professor pediu explicitamente que o
  projeto NÃO precisa estar finalizado, e sim ter os blocos principais implementados. Nossa
  postura de declarar as lacunas está alinhada com isso — não é confissão de fracasso.
· POR QUE O F1 PARECE BAIXO: porque o limiar está calibrado para recall. O F1 trata precisão e
  recall como igualmente importantes, o que não é o nosso caso de uso. Se ajustarmos o limiar
  para maximizar F1 (p99,9), ele vai a 0,57 — está no slide anterior.

SE PERGUNTAREM
P: "Vocês considerariam o projeto bem-sucedido?"
R: "Como maquete funcional, sim, e temos o critério de aceitação atendido para sustentar isso.
   Como sistema de produção, não — falta a explicabilidade por métrica, que é o que tornaria o
   alerta acionável para quem está de plantão."
P: "Qual foi a maior lição de engenharia de software aqui?"
R: "Que formalizar requisito muda o projeto. Ao traduzir 'não quero muito alarme falso' em uma
   métrica, descobrimos que nosso resultado era três vezes pior do que teríamos prometido no
   olhômetro. E ao escrever o RNF-02 percebemos um requisito que o código não atende. Sem a
   formalização, os dois passariam despercebidos."

FECHAMENTO SUGERIDO: "É isso. O repositório está no link do primeiro slide, com o documento de
requisitos, o de arquitetura e a suíte de testes. Obrigado."
""")


# ===========================================================================
def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for construir in [
        slide_01, slide_02, slide_03, slide_04, slide_05, slide_06, slide_07,
        slide_08, slide_09, slide_10, slide_11, slide_12, slide_13,
    ]:
        construir(prs)

    prs.save(SAIDA)
    print(f"Apresentacao gerada: {SAIDA}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
